import os
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Theme Colors
NAVY = RGBColor(15, 32, 67)
BLUE_ACCENT = RGBColor(0, 114, 206)
DARK_GRAY = RGBColor(40, 40, 40)
LIGHT_BG = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)
GRAY_TEXT = RGBColor(90, 90, 90)
LIGHT_BLUE = RGBColor(230, 240, 250)
CARD_BG = RGBColor(255, 255, 255)
BORDER_COLOR = RGBColor(210, 220, 230)
ACCENT_GREEN = RGBColor(34, 139, 34)
ACCENT_ORANGE = RGBColor(218, 112, 214)
ALERT_BG = RGBColor(255, 243, 205)
ALERT_BORDER = RGBColor(255, 222, 117)

def add_header(slide, title_text, category_text="INF-117 — INFORMÁTICA APLICADA A AERONÁUTICA"):
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()
    
    tf = header.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.12)
    
    p0 = tf.paragraphs[0]
    p0.text = category_text.upper()
    p0.font.size = Pt(9)
    p0.font.bold = True
    p0.font.color.rgb = BLUE_ACCENT
    p0.font.name = "Arial"
    
    p1 = tf.add_paragraph()
    p1.text = title_text
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.font.name = "Arial"

def add_footer(slide, current_slide, total_slides=11):
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.25), Inches(9.0), Inches(0.3))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Fatec Sorocaba — CST em Manutenção de Aeronaves  |  Slide {current_slide} de {total_slides}"
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY_TEXT
    p.font.name = "Arial"

def build_aula3():
    prs = pptx.Presentation()
    prs.slide_width = Inches(10.0)
    prs.slide_height = Inches(5.625)
    blank_layout = prs.slide_layouts[6]
    total_slides = 11

    # ==================== SLIDE 1: CAPA ====================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = NAVY
    bg1.line.fill.background()
    
    tb1 = slide1.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8.4), Inches(4.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "FATEC SOROCABA — CURSO SUPERIOR DE TECNOLOGIA EM MANUTENÇÃO DE AERONAVES"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.font.name = "Arial"
    
    p = tf1.add_paragraph()
    p.text = "INF-117 — Informática Aplicada a Aeronáutica"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(200, 225, 255)
    p.font.name = "Arial"
    p.space_before = Pt(8)
    
    p = tf1.add_paragraph()
    p.text = "MS Word Técnico — Normas ABNT (Parte 2)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"
    p.space_before = Pt(14)
    
    p = tf1.add_paragraph()
    p.text = "Quebras de Seção, Numeração ABNT Diferenciada, Figuras/Tabelas e Sumário Automático"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(220, 235, 250)
    p.font.name = "Arial"
    p.space_before = Pt(8)
    
    p = tf1.add_paragraph()
    p.text = "Prof. André Souza  |  Aula 03"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.font.name = "Arial"
    p.space_before = Pt(26)

    # ==================== SLIDE 2: OBJETIVOS ====================
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "Objetivos de Aprendizagem da Aula 03")
    add_footer(slide2, 2, total_slides)
    
    tb2 = slide2.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(4.0))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    
    objs = [
        ("Quebras de Seção vs. Quebras de Página:", " Compreender a diferença estrutural e computacional entre mudar de página e criar partições independentes de cabeçalho."),
        ("Numeração de Páginas ABNT:", " Configurar a contagem desde a Capa, mas exibindo o número impresso apenas a partir da Introdução, no canto superior direito."),
        ("Desvinculação de Cabeçalhos e Rodapés:", " Operar com precisão o comando 'Vincular ao Anterior' para isolar a Capa e o Sumário."),
        ("Tabelas Técnicas e Legendas de Figuras:", " Inserir legendas automáticas numeradas (Figura 1, Tabela 1) e padronizar a indicação de fonte consultada."),
        ("Sumário Automático Dinâmico:", " Gerar o sumário automático com 1 clique a partir dos Estilos e atualizá-lo via atalho F9."),
        ("Exportação Oficial PDF/A:", " Gerar o documento final preservando a conformidade para arquivo de longo prazo (ISO 19005).")
    ]
    
    for i, (title, desc) in enumerate(objs):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_before = Pt(6) if i > 0 else Pt(0)
        
        r1 = p.add_run()
        r1.text = f"•  {title}"
        r1.font.bold = True
        r1.font.size = Pt(11.5)
        r1.font.color.rgb = NAVY
        r1.font.name = "Arial"
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = DARK_GRAY
        r2.font.name = "Arial"

    # ==================== SLIDE 3: O DESAFIO DA ABNT ====================
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "O Desafio da Numeração ABNT (NBR 14724)")
    add_footer(slide3, 3, total_slides)
    
    # 2 Big Cards Layout
    cards_s3 = [
        ("REGRA NORMATIVA DA ABNT", "O que a norma exige?", [
            "Todas as folhas são CONTADAS a partir da Capa (Capa = 1, Sumário = 2, Introdução = 3...).",
            "A numeração só deve aparecer IMPRESSA a partir da primeira folha da parte textual (Introdução).",
            "O número deve ser posicionado no canto superior direito, a 2 cm da borda da folha.",
            "Capa, Sumário e elementos pré-textuais NÃO podem exibir número de página visível!"
        ], NAVY, LIGHT_BLUE),
        ("A SOLUÇÃO COMPUTACIONAL", "Como o MS Word resolve?", [
            "Um documento comum do Word replica o mesmo cabeçalho em todas as páginas.",
            "Para criar regras diferentes, dividimos o arquivo em SEÇÕES independentes.",
            "Seção 1: Capa e Sumário (Cabeçalho em branco).",
            "Seção 2: Introdução e Desenvolvimento (Cabeçalho com número visível iniciando em 3).",
            "A chave de tudo: Desativar o botão 'Vincular ao Anterior'!"
        ], BLUE_ACCENT, WHITE)
    ]
    
    for i, (ctitle, csub, clist, tcolor, bgcolor) in enumerate(cards_s3):
        left_pos = Inches(0.6 + i * 4.5)
        card_bg = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.15), Inches(4.3), Inches(3.9))
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = bgcolor
        card_bg.line.color.rgb = BORDER_COLOR
        
        tb = slide3.shapes.add_textbox(left_pos + Inches(0.2), Inches(1.25), Inches(3.9), Inches(3.7))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = ctitle
        p.font.bold = True
        p.font.size = Pt(12.5)
        p.font.color.rgb = tcolor
        
        p_sub = tf.add_paragraph()
        p_sub.text = csub
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = DARK_GRAY
        p_sub.space_before = Pt(2)
        
        for item in clist:
            pi = tf.add_paragraph()
            pi.text = f"• {item}"
            pi.font.size = Pt(10.5)
            pi.font.color.rgb = DARK_GRAY
            pi.space_before = Pt(6)

    # ==================== SLIDE 4: QUEBRAS DE SEÇÃO VS PAGINA ====================
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "Quebra de Página vs. Quebra de Seção")
    add_footer(slide4, 4, total_slides)
    
    cards_s4 = [
        ("QUEBRA DE PÁGINA COMUM", "Atalho: Ctrl + Enter", [
            "Apenas empurra o cursor para o topo da folha seguinte.",
            "Mantém rigorosamente o mesmo cabeçalho, rodapé e orientação da página anterior.",
            "Se você inserir um número de página após uma quebra simples, ele aparecerá em TODAS as folhas.",
            "Uso: Mudar de página dentro do mesmo capítulo."
        ]),
        ("QUEBRA DE SEÇÃO (PRÓXIMA PÁGINA)", "Menu: Layout -> Quebras -> Seção", [
            "Cria uma barreira computacional isolando o documento em blocos autônomos.",
            "Permite cabeçalhos independentes, rodapés com formatos distintos e até mudar a orientação (Retrato / Paisagem).",
            "Permite desvincular a numeração da Introdução para não poluir a Capa e o Sumário.",
            "Uso: Divisão entre Capa/Sumário (Seção 1) e Conteúdo (Seção 2)."
        ])
    ]
    
    for i, (ctitle, csub, clist) in enumerate(cards_s4):
        left_pos = Inches(0.6 + i * 4.5)
        card_bg = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.15), Inches(4.3), Inches(3.9))
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = CARD_BG if i == 0 else LIGHT_BLUE
        card_bg.line.color.rgb = BORDER_COLOR
        
        tb = slide4.shapes.add_textbox(left_pos + Inches(0.2), Inches(1.25), Inches(3.9), Inches(3.7))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = ctitle
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY
        
        p_sub = tf.add_paragraph()
        p_sub.text = csub
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = BLUE_ACCENT
        p_sub.space_before = Pt(2)
        
        for item in clist:
            pi = tf.add_paragraph()
            pi.text = f"• {item}"
            pi.font.size = Pt(10.5)
            pi.font.color.rgb = DARK_GRAY
            pi.space_before = Pt(6)

    # ==================== SLIDE 5: DESVINCULANDO CABEÇALHOS ====================
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "Passo a Passo: Desvinculando Cabeçalhos e Rodapés")
    add_footer(slide5, 5, total_slides)
    
    tb5 = slide5.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(4.0))
    tf5 = tb5.text_frame
    tf5.word_wrap = True
    
    steps5 = [
        ("Passo 1: Inserir a Quebra de Seção", "Posicione o cursor no final da página do Sumário (Página 2). Acesse Layout -> Quebras -> Quebras de Seção: Próxima Página."),
        ("Passo 2: Ativar o Caractere Invisível (¶)", "Na guia Página Inicial, clique no botão ¶ (Mostrar Tudo) para confirmar visualmente onde a quebra de seção está posicionada."),
        ("Passo 3: Abrir o Cabeçalho da Seção 2", "Vá até a página da Introdução (Seção 2) e dê um DUPLO CLIQUE na área superior da folha (cabeçalho)."),
        ("Passo 4: Desativar 'Vincular ao Anterior'", "Na barra superior 'Cabeçalho e Rodapé', localize o botão 'Vincular ao Anterior' (que estará sombreado/ativo). CLIQUE NELE PARA DESATIVÁ-LO!"),
        ("Resultado:", "Agora o cabeçalho da Seção 2 está 100% independente do cabeçalho da Capa e do Sumário!")
    ]
    
    for i, (title, desc) in enumerate(steps5):
        p = tf5.paragraphs[0] if i == 0 else tf5.add_paragraph()
        p.space_before = Pt(8) if i > 0 else Pt(0)
        
        r1 = p.add_run()
        r1.text = f"{title}: "
        r1.font.bold = True
        r1.font.size = Pt(11.5)
        r1.font.color.rgb = NAVY if i < 4 else RGBColor(0, 140, 60)
        r1.font.name = "Arial"
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = DARK_GRAY
        r2.font.name = "Arial"

    # ==================== SLIDE 6: CONFIGURANDO A NUMERAÇÃO ====================
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "Inserindo e Formatando o Número de Página ABNT")
    add_footer(slide6, 6, total_slides)
    
    # 3 Cards Steps
    cards_s6 = [
        ("1. FORMATAR NÚMERO", "Definir o valor inicial", [
            "Com o cabeçalho da Seção 2 aberto:",
            "Acesse: Número de Página -> Formatar Números de Página...",
            "Marque: 'Iniciar em:' e digite o número 3 (contando Capa=1 e Sumário=2).",
            "Clique em OK."
        ]),
        ("2. POSICIONAR NÚMERO", "Canto superior direito", [
            "Acesse: Número de Página -> Início da Página.",
            "Escolha: 'Número Sem Formatação 3' (alinhado à direita).",
            "O número 3 aparecerá automaticamente na Introdução.",
            "Dê duplo clique no corpo do texto para sair do cabeçalho."
        ]),
        ("3. CONFERÊNCIA", "Verificação obrigatória", [
            "Role o documento até a Capa (Página 1) -> Não deve ter número!",
            "Verifique o Sumário (Página 2) -> Não deve ter número!",
            "Verifique a Introdução (Página 3) -> Exibe o número 3 no canto superior direito!",
            "As páginas seguintes continuam em 4, 5, 6..."
        ])
    ]
    
    for i, (ctitle, csub, clist) in enumerate(cards_s6):
        left_pos = Inches(0.6 + i * 3.0)
        card_bg = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.2), Inches(2.8), Inches(3.8))
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = LIGHT_BLUE if i == 1 else CARD_BG
        card_bg.line.color.rgb = BORDER_COLOR
        
        tb = slide6.shapes.add_textbox(left_pos + Inches(0.15), Inches(1.3), Inches(2.5), Inches(3.6))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = ctitle
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY
        
        p_sub = tf.add_paragraph()
        p_sub.text = csub
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = BLUE_ACCENT
        p_sub.space_before = Pt(2)
        
        for item in clist:
            pi = tf.add_paragraph()
            pi.text = f"• {item}"
            pi.font.size = Pt(10)
            pi.font.color.rgb = DARK_GRAY
            pi.space_before = Pt(6)

    # ==================== SLIDE 7: FIGURAS E TABELAS COM LEGENDAS ====================
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "Tabelas Técnicas e Legendas Automáticas de Figuras")
    add_footer(slide7, 7, total_slides)
    
    cards_s7 = [
        ("LEGENDAS AUTOMÁTICAS DE FIGURAS", "Padrão ABNT para ilustrações", [
            "Clique com o botão direito sobre a imagem -> 'Inserir Legenda...'",
            "Rótulo: selecione 'Figura'.",
            "Digite o título descritivo (ex: ': Aeronave Embraer E195-E2').",
            "Posição: Acima ou abaixo da figura (conforme padrão institucional).",
            "Abaixo da imagem, insira em fonte 10 pt: 'Fonte: ANAC (2026)'."
        ]),
        ("TABELAS TÉCNICAS PROFISSIONAIS", "Estrutura e dados de engenharia", [
            "Guia Inserir -> Tabela -> Definir colunas e linhas.",
            "Formatar cabeçalho com fundo escuro e texto em negrito.",
            "Alinhar dados numéricos e técnicos verticalmente ao centro.",
            "Inserir legenda oficial: botão direito -> Inserir Legenda -> Rótulo: 'Tabela'.",
            "Adicionar 'Fonte: Autores (2026)' no rodapé da tabela em 10 pt."
        ])
    ]
    
    for i, (ctitle, csub, clist) in enumerate(cards_s7):
        left_pos = Inches(0.6 + i * 4.5)
        card_bg = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.15), Inches(4.3), Inches(3.9))
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = WHITE if i == 0 else LIGHT_BLUE
        card_bg.line.color.rgb = BORDER_COLOR
        
        tb = slide7.shapes.add_textbox(left_pos + Inches(0.2), Inches(1.25), Inches(3.9), Inches(3.7))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = ctitle
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY
        
        p_sub = tf.add_paragraph()
        p_sub.text = csub
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = BLUE_ACCENT
        p_sub.space_before = Pt(2)
        
        for item in clist:
            pi = tf.add_paragraph()
            pi.text = f"• {item}"
            pi.font.size = Pt(10.5)
            pi.font.color.rgb = DARK_GRAY
            pi.space_before = Pt(6)

    # ==================== SLIDE 8: SUMÁRIO AUTOMÁTICO ====================
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "Sumário Automático e a Tecla F9")
    add_footer(slide8, 8, total_slides)
    
    tb8 = slide8.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(4.0))
    tf8 = tb8.text_frame
    tf8.word_wrap = True
    
    topics8 = [
        ("Como o Word constrói o Sumário?", " O Word varre todo o documento procurando parágrafos formatados com os estilos Título 1, Título 2 e Título 3, associando o texto com o número exato da página onde eles se encontram."),
        ("Inserção na Página 2:", " Posicione o cursor na página do Sumário -> Acesse a guia Referências -> Sumário -> Sumário Automático 1."),
        ("O que fazer quando o texto for alterado?", " Se você adicionar parágrafos ou mudar títulos, o sumário NÃO muda sozinho na mesma hora. Você deve atualizá-lo!"),
        ("Atalho Dinâmico F9:", " Clique sobre qualquer parte do sumário e pressione a tecla F9 (ou clique com o botão direito -> 'Atualizar Campo') -> Escolha 'Atualizar o índice inteiro' -> OK!"),
        ("Vantagem de Engenharia:", " Elimina para sempre o erro clássico de sumários digitados à mão com pontinhos que não alinham!")
    ]
    
    for i, (title, desc) in enumerate(topics8):
        p = tf8.paragraphs[0] if i == 0 else tf8.add_paragraph()
        p.space_before = Pt(8) if i > 0 else Pt(0)
        
        r1 = p.add_run()
        r1.text = f"•  {title} "
        r1.font.bold = True
        r1.font.size = Pt(11.5)
        r1.font.color.rgb = NAVY
        r1.font.name = "Arial"
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = DARK_GRAY
        r2.font.name = "Arial"

    # ==================== SLIDE 9: EXPORTAÇÃO PDF/A ====================
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "Exportação Oficial no Padrão PDF/A (ISO 19005)")
    add_footer(slide9, 9, total_slides)
    
    cards_s9 = [
        ("POR QUE USAR PDF/A?", "Padrão oficial de arquivamento", [
            "PDF/A é uma versão padronizada pela ISO para arquivamento digital de longo prazo.",
            "Garante que fontes, diagramas e fórmulas fiquem 100% embutidos no arquivo.",
            "Garante que o avaliador ou órgão regulador verá exatamente a mesma diagramação, sem desconfigurações.",
            "Formato obrigatório para entrega de TCC e relatórios técnicos da Fatec."
        ]),
        ("COMO EXPORTAR NO MS WORD", "3 passos simples", [
            "1. Acesse: Arquivo -> Salvar Como (ou Exportar).",
            "2. Selecione o formato: PDF (*.pdf).",
            "3. Clique em 'Mais opções...' (ou 'Opções').",
            "4. Marque a caixa de seleção: 'Compatível com ISO 19005-1 (PDF/A)'.",
            "5. Clique em Salvar."
        ])
    ]
    
    for i, (ctitle, csub, clist) in enumerate(cards_s9):
        left_pos = Inches(0.6 + i * 4.5)
        card_bg = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.15), Inches(4.3), Inches(3.9))
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = CARD_BG if i == 0 else LIGHT_BLUE
        card_bg.line.color.rgb = BORDER_COLOR
        
        tb = slide9.shapes.add_textbox(left_pos + Inches(0.2), Inches(1.25), Inches(3.9), Inches(3.7))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = ctitle
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY
        
        p_sub = tf.add_paragraph()
        p_sub.text = csub
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = BLUE_ACCENT
        p_sub.space_before = Pt(2)
        
        for item in clist:
            pi = tf.add_paragraph()
            pi.text = f"• {item}"
            pi.font.size = Pt(10.5)
            pi.font.color.rgb = DARK_GRAY
            pi.space_before = Pt(6)

    # ==================== SLIDE 10: ATIVIDADE DE FIXAÇÃO ====================
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "Atividade Prática de Fixação no Laboratório")
    add_footer(slide10, 10, total_slides)
    
    tb10 = slide10.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(4.0))
    tf10 = tb10.text_frame
    tf10.word_wrap = True
    
    tasks10 = [
        ("1. Abrir o Documento da Aula 2:", " Abra o relatório com os estilos ABNT já configurados na aula passada."),
        ("2. Inserir Quebra de Seção:", " Crie a fronteira de seção no final do Sumário."),
        ("3. Desvincular e Numerar:", " Desative 'Vincular ao Anterior', configure a numeração iniciando em 3 e posicione no topo direito."),
        ("4. Legendas e Tabela:", " Adicione a tabela técnica com cabeçalho formatado e a figura com legenda oficial."),
        ("5. Sumário Automático & PDF/A:", " Insira o Sumário Automático 1 e exporte o arquivo final no formato PDF/A.")
    ]
    
    for i, (title, desc) in enumerate(tasks10):
        p = tf10.paragraphs[0] if i == 0 else tf10.add_paragraph()
        p.space_before = Pt(8) if i > 0 else Pt(0)
        
        r1 = p.add_run()
        r1.text = f"{title} "
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = NAVY
        r1.font.name = "Arial"
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(12)
        r2.font.color.rgb = DARK_GRAY
        r2.font.name = "Arial"

    # ==================== SLIDE 11: ATALHOS NO WORD ====================
    slide11 = prs.slides.add_slide(blank_layout)
    add_header(slide11, "Dicas de Produtividade & Atalhos no MS Word")
    add_footer(slide11, 11, total_slides)
    
    shortcuts = [
        ("Ctrl + Shift + 8 (Ctrl + *)", "Exibe ou oculta os caracteres não imprimíveis (¶) e quebras de seção"),
        ("Alt + Shift + P", "Insere o campo dinâmico de Número de Página"),
        ("F9", "Atualiza campos selecionados (como o Sumário Automático)"),
        ("Ctrl + Enter", "Insere uma Quebra de Página simples"),
        ("F12", "Abre a janela 'Salvar Como' instantaneamente"),
        ("Ctrl + J", "Aplica o alinhamento Justificado ao parágrafo"),
        ("Ctrl + E", "Aplica o alinhamento Centralizado ao parágrafo")
    ]
    
    # Table layout for shortcuts
    table_shape = slide11.shapes.add_table(8, 2, Inches(0.8), Inches(1.2), Inches(8.4), Inches(3.6))
    table = table_shape.table
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(5.2)
    
    # Table header
    cell_00 = table.cell(0, 0)
    cell_00.text = "ATALHO DE TECLADO"
    cell_00.fill.solid()
    cell_00.fill.fore_color.rgb = NAVY
    cell_00.text_frame.paragraphs[0].font.bold = True
    cell_00.text_frame.paragraphs[0].font.color.rgb = WHITE
    cell_00.text_frame.paragraphs[0].font.size = Pt(11)
    
    cell_01 = table.cell(0, 1)
    cell_01.text = "FUNÇÃO NO MICROSOFT WORD"
    cell_01.fill.solid()
    cell_01.fill.fore_color.rgb = NAVY
    cell_01.text_frame.paragraphs[0].font.bold = True
    cell_01.text_frame.paragraphs[0].font.color.rgb = WHITE
    cell_01.text_frame.paragraphs[0].font.size = Pt(11)
    
    for row_idx, (sc, desc) in enumerate(shortcuts, start=1):
        c0 = table.cell(row_idx, 0)
        c0.text = sc
        c0.fill.solid()
        c0.fill.fore_color.rgb = LIGHT_BLUE if row_idx % 2 == 1 else WHITE
        p = c0.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(10.5)
        p.font.color.rgb = NAVY
        
        c1 = table.cell(row_idx, 1)
        c1.text = desc
        c1.fill.solid()
        c1.fill.fore_color.rgb = LIGHT_BLUE if row_idx % 2 == 1 else WHITE
        p = c1.text_frame.paragraphs[0]
        p.font.size = Pt(10.5)
        p.font.color.rgb = DARK_GRAY

    output_path = r'c:\projetos\Material\Informática Aplicada a Aeronáutica\Slides\Slides_Aula_03_Word_Tecnico_ABNT_Parte2.pptx'
    prs.save(output_path)
    print(f"Presentation saved successfully: {output_path}")


def build_aula4():
    prs = pptx.Presentation()
    prs.slide_width = Inches(10.0)
    prs.slide_height = Inches(5.625)
    blank_layout = prs.slide_layouts[6]
    total_slides = 11

    # ==================== SLIDE 1: CAPA ====================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = NAVY
    bg1.line.fill.background()
    
    tb1 = slide1.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8.4), Inches(4.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "FATEC SOROCABA — CURSO SUPERIOR DE TECNOLOGIA EM MANUTENÇÃO DE AERONAVES"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.font.name = "Arial"
    
    p = tf1.add_paragraph()
    p.text = "INF-117 — Informática Aplicada a Aeronáutica"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(200, 225, 255)
    p.font.name = "Arial"
    p.space_before = Pt(8)
    
    p = tf1.add_paragraph()
    p.text = "Oficina Prática & Exercício Avaliativo de Formatação ABNT"
    p.font.size = Pt(23)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"
    p.space_before = Pt(14)
    
    p = tf1.add_paragraph()
    p.text = "Laboratório de Execução Orientada, Checklist de Entrega e Exercício Avaliativo 1"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(220, 235, 250)
    p.font.name = "Arial"
    p.space_before = Pt(8)
    
    p = tf1.add_paragraph()
    p.text = "Prof. André Souza  |  Aula 04"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.font.name = "Arial"
    p.space_before = Pt(26)

    # ==================== SLIDE 2: OBJETIVOS ====================
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "Objetivos da Oficina Prática de Formatação ABNT")
    add_footer(slide2, 2, total_slides)
    
    tb2 = slide2.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(4.0))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    
    objs4 = [
        ("Consolidação Prática Individual:", " Aplicar de forma autônoma e passo a passo todas as regras de formatação ABNT trabalhadas nas Aulas 2 e 3."),
        ("Transformação de Texto Bruto:", " Receber o arquivo de texto não-formatado sobre Introdução à Aviação e estruturá-lo como um Relatório Técnico oficial de engenharia."),
        ("Atendimento Bancada a Bancada:", " Esclarecer dúvidas pontuais com o professor durante o desenvolvimento em laboratório."),
        ("Conclusão do Exercício Avaliativo 1:", " Concluir e entregar o trabalho para somar pontuação na média de exercícios práticos contínuos (peso de 50%)."),
        ("Entrega em Múltiplos Formatos:", " Gerar o arquivo editável (.docx) e o arquivo final inviolável no padrão oficial PDF/A (.pdf).")
    ]
    
    for i, (title, desc) in enumerate(objs4):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_before = Pt(8) if i > 0 else Pt(0)
        
        r1 = p.add_run()
        r1.text = f"•  {title} "
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = NAVY
        r1.font.name = "Arial"
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(12)
        r2.font.color.rgb = DARK_GRAY
        r2.font.name = "Arial"

    # ==================== SLIDE 3: DINÂMICA DA AULA ====================
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "Dinâmica e Gestão do Tempo no Laboratório (~1h40)")
    add_footer(slide3, 3, total_slides)
    
    cards_s3_4 = [
        ("1. SETUP & ALINHAMENTO", "Primeiros 15 minutos", [
            "Ligação dos computadores e login no Office 365.",
            "Download do texto base: Texto_Exercicio_Historia_da_Aviacao.txt.",
            "Abertura do Guia_Exercicio_Formatacao_Word.md.",
            "Revisão rápida dos critérios de pontuação pelo professor."
        ]),
        ("2. DESENVOLVIMENTO", "60 minutos de prática", [
            "Execução autônoma do checklist de formatação pelo estudante.",
            "Configuração de margens, estilos, quebras de seção e numeração.",
            "Formatação da tabela de aeronaves e legenda de figura.",
            "Atendimento individual do professor passando nas bancadas."
        ]),
        ("3. CONFERÊNCIA & ENVIO", "Últimos 15 minutos", [
            "Geração e atualização final do Sumário Automático.",
            "Conferência do checklist completo de requisitos ABNT.",
            "Exportação no formato PDF/A.",
            "Upload dos 2 arquivos (.docx e .pdf) na plataforma institucional."
        ])
    ]
    
    for i, (ctitle, csub, clist) in enumerate(cards_s3_4):
        left_pos = Inches(0.6 + i * 3.0)
        card_bg = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.2), Inches(2.8), Inches(3.8))
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = LIGHT_BLUE if i == 1 else CARD_BG
        card_bg.line.color.rgb = BORDER_COLOR
        
        tb = slide3.shapes.add_textbox(left_pos + Inches(0.15), Inches(1.3), Inches(2.5), Inches(3.6))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = ctitle
        p.font.bold = True
        p.font.size = Pt(11.5)
        p.font.color.rgb = NAVY
        
        p_sub = tf.add_paragraph()
        p_sub.text = csub
        p_sub.font.size = Pt(9.5)
        p_sub.font.color.rgb = BLUE_ACCENT
        p_sub.space_before = Pt(2)
        
        for item in clist:
            pi = tf.add_paragraph()
            pi.text = f"• {item}"
            pi.font.size = Pt(9.5)
            pi.font.color.rgb = DARK_GRAY
            pi.space_before = Pt(5)

    # ==================== SLIDE 4: O ARQUIVO BASE ====================
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "O Arquivo Base — Introdução à Aviação Civil")
    add_footer(slide4, 4, total_slides)
    
    tb4 = slide4.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(4.0))
    tf4 = tb4.text_frame
    tf4.word_wrap = True
    
    sections = [
        ("Seção 1: Santos Dumont e o Voo do 14-Bis", "História, pioneirismo aeronáutico e evolução dos primeiros voos na França."),
        ("Seção 2: A Criação da ANAC e o Papel Regulador", "Histórico do DAC à ANAC, regulamentos RBAC e segurança operacional."),
        ("Seção 3: Aeronaves Leves Esportivas (ALE / LSA)", "Definição técnica de aeronaves experimentais e desportivas leves."),
        ("Seção 4: A Evolução da Aviação Comercial e a Embraer", "Crescimento da aviação comercial a jato e a liderança mundial dos E-Jets da Embraer."),
        ("Seção 5: Formação e Licenças na Aviação Civil", "Requisitos para obtenção de CMA, licenças de piloto e habilitações de mecânicos CHT."),
        ("Seção 6: Segurança Operacional e Tomada de Decisão", "Cultura de segurança de voo e gerenciamento de risco aeronáutico."),
        ("Seção 7: Considerações Finais e Perspectivas", "Sustentabilidade, combustíveis SAF e inovação na manutenção de aeronaves.")
    ]
    
    for i, (title, desc) in enumerate(sections):
        p = tf4.paragraphs[0] if i == 0 else tf4.add_paragraph()
        p.space_before = Pt(5) if i > 0 else Pt(0)
        
        r1 = p.add_run()
        r1.text = f"{title}: "
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = NAVY
        r1.font.name = "Arial"
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(11)
        r2.font.color.rgb = DARK_GRAY
        r2.font.name = "Arial"

    # ==================== SLIDE 5: CHECKLIST 1 - MARGENS & ESTILOS ====================
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "Checklist ABNT 1 — Margens, Papel e Estilos de Texto")
    add_footer(slide5, 5, total_slides)
    
    cards_s5_4 = [
        ("CONFIGURAÇÃO DE PÁGINA", "Layout -> Margens", [
            "Tamanho do Papel: A4 (21,0 cm x 29,7 cm).",
            "Margem Superior: 3,0 cm.",
            "Margem Esquerda: 3,0 cm.",
            "Margem Inferior: 2,0 cm.",
            "Margem Direita: 2,0 cm."
        ]),
        ("ESTILO NORMAL (CORPO DE TEXTO)", "Página Inicial -> Estilos", [
            "Fonte: Arial ou Calibri 12 pt.",
            "Alinhamento: Justificado (Ctrl + J).",
            "Espaçamento entre Linhas: 1,5 linha.",
            "Recuo Especial de 1ª Linha: 1,25 cm.",
            "Espaçamento Depois: 6 pt."
        ]),
        ("ESTILOS DE TÍTULOS", "Hierarquia estrutural", [
            "Título 1: Arial/Calibri 14 pt, Negrito, MAIÚSCULAS (Antes 12 pt / Depois 6 pt).",
            "Título 2: Arial/Calibri 12 pt, Negrito, Caixa Mista (Antes 6 pt / Depois 6 pt).",
            "Aplicar nas seções principais 1 a 7 e subseções."
        ])
    ]
    
    for i, (ctitle, csub, clist) in enumerate(cards_s5_4):
        left_pos = Inches(0.6 + i * 3.0)
        card_bg = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.2), Inches(2.8), Inches(3.8))
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = CARD_BG if i == 0 else (LIGHT_BLUE if i == 1 else CARD_BG)
        card_bg.line.color.rgb = BORDER_COLOR
        
        tb = slide5.shapes.add_textbox(left_pos + Inches(0.15), Inches(1.3), Inches(2.5), Inches(3.6))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = ctitle
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = NAVY
        
        p_sub = tf.add_paragraph()
        p_sub.text = csub
        p_sub.font.size = Pt(9.5)
        p_sub.font.color.rgb = BLUE_ACCENT
        p_sub.space_before = Pt(2)
        
        for item in clist:
            pi = tf.add_paragraph()
            pi.text = f"• {item}"
            pi.font.size = Pt(9.5)
            pi.font.color.rgb = DARK_GRAY
            pi.space_before = Pt(5)

    # ==================== SLIDE 6: CHECKLIST 2 - SEÇÕES & NUMERAÇÃO ====================
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "Checklist ABNT 2 — Quebras de Seção e Numeração")
    add_footer(slide6, 6, total_slides)
    
    tb6 = slide6.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(4.0))
    tf6 = tb6.text_frame
    tf6.word_wrap = True
    
    chk6 = [
        ("Página 1 (Capa):", " Centralizar dados institucionais (Fatec Sorocaba, INF-117, título do relatório, autor, local e ano). Inserir Quebra de Seção (Próxima Página)."),
        ("Página 2 (Sumário):", " Reservar para o Sumário Automático. Inserir Quebra de Seção (Próxima Página) no final da folha."),
        ("Página 3 (Introdução - Seção 2):", " Dê duplo clique no cabeçalho e DESATIVE o botão 'Vincular ao Anterior'."),
        ("Inserção do Número de Página:", " Número de Página -> Formatar Números de Página -> Marcar 'Iniciar em: 3' -> Número de Página -> Início da Página -> Número Sem Formatação 3 (canto direito)."),
        ("Conferência Visual Obrigatória:", " Capa e Sumário SEM NÚMERO VISÍVEL; numeração iniciando visível em 3 na Introdução!")
    ]
    
    for i, (title, desc) in enumerate(chk6):
        p = tf6.paragraphs[0] if i == 0 else tf6.add_paragraph()
        p.space_before = Pt(8) if i > 0 else Pt(0)
        
        r1 = p.add_run()
        r1.text = f"•  {title} "
        r1.font.bold = True
        r1.font.size = Pt(11.5)
        r1.font.color.rgb = NAVY
        r1.font.name = "Arial"
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = DARK_GRAY
        r2.font.name = "Arial"

    # ==================== SLIDE 7: CHECKLIST 3 - TABELA & FIGURAS ====================
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "Checklist ABNT 3 — Tabela de Aeronaves e Figuras")
    add_footer(slide7, 7, total_slides)
    
    cards_s7_4 = [
        ("TABELA DE AERONAVES (SEÇÃO 4)", "Grade estruturada com 7 colunas x 8 linhas", [
            "Inserir tabela com dados de Santos Dumont 14-Bis, Demoiselle, DC-3, Boeing 707, Concorde e E195-E2.",
            "Cabeçalho formatado com fundo escuro e texto em negrito.",
            "Propriedades da Tabela -> Habilitar 'Repetir Linhas de Cabeçalho'.",
            "Inserir legenda acima: 'Tabela 1: Comparativo de especificações de aeronaves'.",
            "Inserir fonte abaixo em 10 pt: 'Fonte: Cartilha ANAC e Manuais (2026)'."
        ]),
        ("FIGURA TÉCNICA COM LEGENDA", "Ilustração oficial no documento", [
            "Inserir imagem técnica centralizada (Embraer E195-E2 ou 14-Bis).",
            "Botão direito -> 'Inserir Legenda...'.",
            "Rótulo: 'Figura' -> ': Aeronave comercial com motores turbofan'.",
            "Inserir linha de fonte abaixo da imagem em 10 pt: 'Fonte: Embraer Media (2026)'."
        ])
    ]
    
    for i, (ctitle, csub, clist) in enumerate(cards_s7_4):
        left_pos = Inches(0.6 + i * 4.5)
        card_bg = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.15), Inches(4.3), Inches(3.9))
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = LIGHT_BLUE if i == 0 else WHITE
        card_bg.line.color.rgb = BORDER_COLOR
        
        tb = slide7.shapes.add_textbox(left_pos + Inches(0.2), Inches(1.25), Inches(3.9), Inches(3.7))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = ctitle
        p.font.bold = True
        p.font.size = Pt(11.5)
        p.font.color.rgb = NAVY
        
        p_sub = tf.add_paragraph()
        p_sub.text = csub
        p_sub.font.size = Pt(9.5)
        p_sub.font.color.rgb = BLUE_ACCENT
        p_sub.space_before = Pt(2)
        
        for item in clist:
            pi = tf.add_paragraph()
            pi.text = f"• {item}"
            pi.font.size = Pt(10)
            pi.font.color.rgb = DARK_GRAY
            pi.space_before = Pt(6)

    # ==================== SLIDE 8: CHECKLIST 4 - SUMÁRIO & PDF/A ====================
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "Checklist ABNT 4 — Sumário Automático e PDF/A")
    add_footer(slide8, 8, total_slides)
    
    cards_s8_4 = [
        ("GERAÇÃO DO SUMÁRIO AUTOMÁTICO", "Página 2 reservada", [
            "Posicionar o cursor na Página 2 reservada.",
            "Guia Referências -> Sumário -> Sumário Automático 1.",
            "Verificar se as 7 seções foram listadas com os números de página exatos.",
            "Pressionar F9 para atualizar o índice antes de exportar."
        ]),
        ("EXPORTAÇÃO OFICIAL PDF/A", "ISO 19005 para arquivamento", [
            "Arquivo -> Salvar Como -> Tipo: PDF (*.pdf).",
            "Mais opções... -> Marcar 'Compatível com ISO 19005-1 (PDF/A)'.",
            "Salvar com nome padronizado.",
            "Abrir o PDF para conferência visual das margens e numeração."
        ])
    ]
    
    for i, (ctitle, csub, clist) in enumerate(cards_s8_4):
        left_pos = Inches(0.6 + i * 4.5)
        card_bg = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.15), Inches(4.3), Inches(3.9))
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = CARD_BG if i == 0 else LIGHT_BLUE
        card_bg.line.color.rgb = BORDER_COLOR
        
        tb = slide8.shapes.add_textbox(left_pos + Inches(0.2), Inches(1.25), Inches(3.9), Inches(3.7))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = ctitle
        p.font.bold = True
        p.font.size = Pt(11.5)
        p.font.color.rgb = NAVY
        
        p_sub = tf.add_paragraph()
        p_sub.text = csub
        p_sub.font.size = Pt(9.5)
        p_sub.font.color.rgb = BLUE_ACCENT
        p_sub.space_before = Pt(2)
        
        for item in clist:
            pi = tf.add_paragraph()
            pi.text = f"• {item}"
            pi.font.size = Pt(10)
            pi.font.color.rgb = DARK_GRAY
            pi.space_before = Pt(6)

    # ==================== SLIDE 9: CRITÉRIOS DE AVALIAÇÃO ====================
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "Critérios de Correção — Exercício Avaliativo 1")
    add_footer(slide9, 9, total_slides)
    
    criteria = [
        ("Configuração de Margens e Papel", "1,0 pt", "Margens corretas (3-3-2-2) e tamanho A4."),
        ("Aplicação de Estilos (Normal, Título 1, Título 2)", "2,5 pts", "Uso real dos estilos, fontes, recuo de 1,25 cm e entrelinhas 1,5."),
        ("Quebras de Seção e Numeração ABNT", "3,0 pts", "Desvinculação correta; numeração visível iniciando em 3 na Introdução."),
        ("Tabelas, Figuras e Legendas Oficiais", "1,5 pts", "Legendas automáticas formatadas e indicação de fonte consultada."),
        ("Sumário Automático e Exportação PDF/A", "2,0 pts", "Sumário dinâmico funcionando e arquivo PDF gerado corretamente."),
        ("TOTAL DO EXERCÍCIO 1", "10,0 pts", "Composição da média de exercícios práticos contínuos (peso 50%).")
    ]
    
    table_shape = slide9.shapes.add_table(7, 3, Inches(0.6), Inches(1.2), Inches(8.8), Inches(3.6))
    table = table_shape.table
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(1.2)
    table.columns[2].width = Inches(4.4)
    
    headers = ["CRITÉRIO AVALIADO", "PONTUAÇÃO", "O QUE SERÁ VERIFICADO"]
    for col_idx, htext in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = htext
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.size = Pt(10.5)
        
    for row_idx, (crit, pts, desc) in enumerate(criteria, start=1):
        c0 = table.cell(row_idx, 0)
        c0.text = crit
        c0.fill.solid()
        c0.fill.fore_color.rgb = LIGHT_BLUE if row_idx % 2 == 1 else WHITE
        p = c0.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = NAVY
        
        c1 = table.cell(row_idx, 1)
        c1.text = pts
        c1.fill.solid()
        c1.fill.fore_color.rgb = LIGHT_BLUE if row_idx % 2 == 1 else WHITE
        p = c1.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = BLUE_ACCENT
        
        c2 = table.cell(row_idx, 2)
        c2.text = desc
        c2.fill.solid()
        c2.fill.fore_color.rgb = LIGHT_BLUE if row_idx % 2 == 1 else WHITE
        p = c2.text_frame.paragraphs[0]
        p.font.size = Pt(9.5)
        p.font.color.rgb = DARK_GRAY

    # ==================== SLIDE 10: NOMENCLATURA E ENVIO ====================
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "Instruções de Envio e Nomenclatura dos Arquivos")
    add_footer(slide10, 10, total_slides)
    
    cards_s10_4 = [
        ("ARQUIVO EDITÁVEL (.DOCX)", "Salvar para arquivo", [
            "Nomenclatura obrigatória:",
            "Relatorio_ABNT_SeuNome_RA.docx",
            "Contém todos os estilos ativos e quebras de seção editáveis."
        ]),
        ("ARQUIVO EXPORTADO (.PDF)", "Padrão PDF/A", [
            "Nomenclatura obrigatória:",
            "Relatorio_ABNT_SeuNome_RA.pdf",
            "Versão final inviolável para avaliação e conferência gráfica."
        ]),
        ("UPLOAD NA PLATAFORMA", "Envio institucional", [
            "Anexar os DOIS arquivos (.docx e .pdf) na tarefa da disciplina.",
            "Conferir se o envio foi concluído com sucesso antes de sair do laboratório."
        ])
    ]
    
    for i, (ctitle, csub, clist) in enumerate(cards_s10_4):
        left_pos = Inches(0.6 + i * 3.0)
        card_bg = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.2), Inches(2.8), Inches(3.8))
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = LIGHT_BLUE if i == 1 else CARD_BG
        card_bg.line.color.rgb = BORDER_COLOR
        
        tb = slide10.shapes.add_textbox(left_pos + Inches(0.15), Inches(1.3), Inches(2.5), Inches(3.6))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = ctitle
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = NAVY
        
        p_sub = tf.add_paragraph()
        p_sub.text = csub
        p_sub.font.size = Pt(9.5)
        p_sub.font.color.rgb = BLUE_ACCENT
        p_sub.space_before = Pt(2)
        
        for item in clist:
            pi = tf.add_paragraph()
            pi.text = f"• {item}"
            pi.font.size = Pt(9.5)
            pi.font.color.rgb = DARK_GRAY
            pi.space_before = Pt(5)

    # ==================== SLIDE 11: PRÓXIMOS PASSOS - EXCEL ====================
    slide11 = prs.slides.add_slide(blank_layout)
    add_header(slide11, "Próximos Passos — Módulo de Microsoft Excel!")
    add_footer(slide11, 11, total_slides)
    
    tb11 = slide11.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(3.8))
    tf11 = tb11.text_frame
    tf11.word_wrap = True
    
    p = tf11.paragraphs[0]
    p.text = "Parabéns pela conclusão do Módulo de Documentação Técnica e ABNT no Word!"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    p = tf11.add_paragraph()
    p.text = "A partir da próxima aula (Aula 05), iniciaremos o Módulo de Microsoft Excel Aplicado à Aeronáutica:"
    p.font.size = Pt(12.5)
    p.font.color.rgb = DARK_GRAY
    p.space_before = Pt(10)
    
    next_topics = [
        "Aula 05: Primeiros passos no Excel (Interface, tipos de células e operadores básicos).",
        "Aula 06: Funções essenciais (SOMA, MÉDIA) e Referências Absolutas ($) com conversão de unidades.",
        "Aula 07: Modelagem de Peso e Balanceamento de Aeronaves (Centro de Gravidade - CG).",
        "Aulas 08 a 13: Lógica condicional (SE), buscas (PROCX), gráficos de engenharia e tabelas dinâmicas.",
        "Avaliação Final: Prova Prática de Excel em 27/11 (50% da média) + Exercícios Contínuos (50%)."
    ]
    
    for item in next_topics:
        pi = tf11.add_paragraph()
        pi.text = f"✈  {item}"
        pi.font.size = Pt(11)
        pi.font.color.rgb = NAVY
        pi.space_before = Pt(8)

    output_path = r'c:\projetos\Material\Informática Aplicada a Aeronáutica\Slides\Slides_Aula_04_Oficina_Pratica_Exercicio_ABNT.pptx'
    prs.save(output_path)
    print(f"Presentation saved successfully: {output_path}")

if __name__ == '__main__':
    build_aula3()
    build_aula4()
