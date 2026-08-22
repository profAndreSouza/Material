import os
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Theme Colors
NAVY = RGBColor(15, 32, 67)
BLUE_ACCENT = RGBColor(0, 114, 206)
DARK_GRAY = RGBColor(40, 40, 40)
LIGHT_BG = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)
GRAY_TEXT = RGBColor(90, 90, 90)
LIGHT_BLUE = RGBColor(230, 240, 250)
CARD_BG = RGBColor(255, 255, 255)
BORDER_COLOR = RGBColor(210, 220, 230)

def add_header(slide, title_text, category_text="INF-117 — INFORMÁTICA APLICADA A AERONÁUTICA"):
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()
    
    tf = header.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.12)
    
    p0 = tf.paragraphs[0]
    p0.text = category_text.upper()
    p0.font.size = Pt(9)
    p0.font.bold = True
    p0.font.color.rgb = BLUE_ACCENT
    p0.font.name = "Arial"
    
    p1 = tf.add_paragraph()
    p1.text = title_text
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.font.name = "Arial"

def add_footer(slide, current_slide, total_slides=10):
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.25), Inches(9.0), Inches(0.3))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Fatec Sorocaba — CST em Manutenção de Aeronaves  |  Slide {current_slide} de {total_slides}"
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY_TEXT
    p.font.name = "Arial"

# ==============================================================================
# AULA 02: MS WORD TÉCNICO - NORMAS ABNT (PARTE 1: MARGENS, TIPOGRAFIA E ESTILOS)
# ==============================================================================
def build_aula2():
    prs = pptx.Presentation()
    prs.slide_width = Inches(10.0)
    prs.slide_height = Inches(5.625)
    blank_layout = prs.slide_layouts[6]
    total_slides = 10

    # Slide 1: Capa
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = NAVY
    bg1.line.fill.background()
    
    tb1 = slide1.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8.4), Inches(4.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "FATEC SOROCABA — CURSO SUPERIOR DE TECNOLOGIA EM MANUTENÇÃO DE AERONAVES"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.font.name = "Arial"
    
    p = tf1.add_paragraph()
    p.text = "INF-117 — Informática Aplicada a Aeronáutica"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(200, 225, 255)
    p.font.name = "Arial"
    p.space_before = Pt(8)
    
    p = tf1.add_paragraph()
    p.text = "MS Word Técnico — Normas ABNT (Parte 1)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"
    p.space_before = Pt(14)
    
    p = tf1.add_paragraph()
    p.text = "Configuração de Página, Margens Oficiais, Tipografia e Criação de Estilos de Parágrafo"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(220, 235, 250)
    p.font.name = "Arial"
    p.space_before = Pt(8)
    
    p = tf1.add_paragraph()
    p.text = "Prof. André Souza  |  Aula 02"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.font.name = "Arial"
    p.space_before = Pt(26)

    # Slide 2: Objetivos
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "Objetivos de Aprendizagem da Aula 02")
    add_footer(slide2, 2, total_slides)
    
    tb2 = slide2.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(4.0))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    
    objs2 = [
        ("Importância Normativa (NBR 14724):", " Compreender a necessidade de padronização visual e estrutural em relatórios de engenharia e aviação."),
        ("Configuração de Página e Margens:", " Definir o papel A4 e as margens oficiais (Superior: 3cm, Esquerda: 3cm, Inferior: 2cm, Direita: 2cm)."),
        ("Parágrafos Técnicos:", " Aplicar alinhamento Justificado, entrelinhas 1,5, recuo de primeira linha em 1,25 cm e espaçamento entre blocos."),
        ("Hierarquia de Estilos de Texto:", " Criar, modificar e aplicar os estilos Normal (corpo), Título 1 (seções) e Título 2 (subseções)."),
        ("Painel de Navegação (Ctrl + F):", " Utilizar o mapa interativo de títulos para inspecionar e navegar no documento."),
        ("Base para o Sumário:", " Preparar o documento para que na Aula 03 seja possível gerar o Sumário Automático instantâneo.")
    ]
    for i, (title, desc) in enumerate(objs2):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_before = Pt(7) if i > 0 else Pt(0)
        r1 = p.add_run()
        r1.text = f"•  {title} "
        r1.font.bold = True
        r1.font.size = Pt(11.5)
        r1.font.color.rgb = NAVY
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = DARK_GRAY

    # Slide 3: Por que usar Estilos?
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "Por que o uso de Estilos é Obrigatório em Engenharia?")
    add_footer(slide3, 3, total_slides)
    
    cards_s3 = [
        ("O PROBLEMA DA FORMATAÇÃO MANUAL", "Vícios comuns que geram retrabalho", [
            "Selecionar parágrafo por parágrafo para trocar fontes e tamanhos gera despadronização.",
            "Qualquer ajuste exige alterar dezenas de páginas manualmente uma a uma.",
            "O Word NÃO reconhece títulos formatados manualmente para gerar o Sumário!",
            "Gera arquivos desorganizados e rejeitados em bancas de TCC e auditorias da ANAC."
        ]),
        ("A SOLUÇÃO ATRAVÉS DE ESTILOS", "Automação e padronização global", [
            "Estilos são regras globais associadas aos blocos do documento.",
            "Alterar o estilo Título 1 atualiza TODOS os capítulos em 1 segundo!",
            "Habilita o Painel de Navegação (Ctrl + F) para saltar entre seções.",
            "Permite ao Word ler a hierarquia e gerar o Sumário Automático com 1 clique."
        ])
    ]
    for i, (ctitle, csub, clist) in enumerate(cards_s3):
        left_pos = Inches(0.6 + i * 4.5)
        card_bg = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.15), Inches(4.3), Inches(3.9))
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = CARD_BG if i == 0 else LIGHT_BLUE
        card_bg.line.color.rgb = BORDER_COLOR
        
        tb = slide3.shapes.add_textbox(left_pos + Inches(0.2), Inches(1.25), Inches(3.9), Inches(3.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ctitle
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY
        p_sub = tf.add_paragraph()
        p_sub.text = csub
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = BLUE_ACCENT
        p_sub.space_before = Pt(2)
        for item in clist:
            pi = tf.add_paragraph()
            pi.text = f"• {item}"
            pi.font.size = Pt(10.5)
            pi.font.color.rgb = DARK_GRAY
            pi.space_before = Pt(6)

    # Slide 4: Padrão Normativo ABNT
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "Padrão Normativo ABNT (NBR 14724) no Word")
    add_footer(slide4, 4, total_slides)
    
    tb4 = slide4.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(4.0))
    tf4 = tb4.text_frame
    tf4.word_wrap = True
    
    rules4 = [
        ("Tamanho do Papel:", " Formato A4 (21,0 cm de largura x 29,7 cm de altura)."),
        ("Margens Superiores e Esquerdas (3,0 cm):", " Espaço maior reservado para encadernação e perfuração oficial do relatório."),
        ("Margens Inferiores e Direitas (2,0 cm):", " Área de respiro e posicionamento da numeração de páginas."),
        ("Tipografia Padrão:", " Fonte Arial ou Calibri ou Times New Roman (tamanho 12 pt no corpo de texto)."),
        ("Espaçamento entre Linhas:", " 1,5 linha em todo o corpo do documento (facilita leitura e anotações do revisor)."),
        ("Recuo de Parágrafo:", " Recuo especial de primeira linha obrigatório em 1,25 cm."),
        ("Alinhamento:", " Sempre Justificado (Ctrl + J), alinhando o texto nas duas margens simultaneamente.")
    ]
    for i, (title, desc) in enumerate(rules4):
        p = tf4.paragraphs[0] if i == 0 else tf4.add_paragraph()
        p.space_before = Pt(6) if i > 0 else Pt(0)
        r1 = p.add_run()
        r1.text = f"•  {title} "
        r1.font.bold = True
        r1.font.size = Pt(11.5)
        r1.font.color.rgb = NAVY
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = DARK_GRAY

    # Slide 5: Hierarquia de Estilos
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "Hierarquia de Estilos: Normal, Título 1 e Título 2")
    add_footer(slide5, 5, total_slides)
    
    cards_s5 = [
        ("ESTILO NORMAL", "Corpo do Relatório", [
            "Fonte: Arial 12 pt.",
            "Alinhamento: Justificado.",
            "Entrelinhas: 1,5 linha.",
            "Recuo 1ª Linha: 1,25 cm.",
            "Espaçamento Depois: 6 pt."
        ]),
        ("ESTILO TÍTULO 1", "Capítulos Principais", [
            "Fonte: Arial 14 pt.",
            "Formato: Negrito, MAIÚSCULAS.",
            "Entrelinhas: 1,5 linha.",
            "Espaçamento Antes: 12 pt.",
            "Espaçamento Depois: 6 pt."
        ]),
        ("ESTILO TÍTULO 2", "Subseções Técnicas", [
            "Fonte: Arial 12 pt.",
            "Formato: Negrito, Caixa Mista.",
            "Entrelinhas: 1,5 linha.",
            "Espaçamento Antes: 6 pt.",
            "Espaçamento Depois: 6 pt."
        ])
    ]
    for i, (ctitle, csub, clist) in enumerate(cards_s5):
        left_pos = Inches(0.6 + i * 3.0)
        card_bg = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.2), Inches(2.8), Inches(3.8))
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = CARD_BG if i == 0 else (LIGHT_BLUE if i == 1 else CARD_BG)
        card_bg.line.color.rgb = BORDER_COLOR
        
        tb = slide5.shapes.add_textbox(left_pos + Inches(0.15), Inches(1.3), Inches(2.5), Inches(3.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ctitle
        p.font.bold = True
        p.font.size = Pt(11.5)
        p.font.color.rgb = NAVY
        p_sub = tf.add_paragraph()
        p_sub.text = csub
        p_sub.font.size = Pt(9.5)
        p_sub.font.color.rgb = BLUE_ACCENT
        p_sub.space_before = Pt(2)
        for item in clist:
            pi = tf.add_paragraph()
            pi.text = f"• {item}"
            pi.font.size = Pt(10)
            pi.font.color.rgb = DARK_GRAY
            pi.space_before = Pt(6)

    # Slide 6: Passo a Passo - Modificando Estilo Normal
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "Passo a Passo — Configurando o Estilo Normal")
    add_footer(slide6, 6, total_slides)
    
    tb6 = slide6.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(4.0))
    tf6 = tb6.text_frame
    tf6.word_wrap = True
    
    steps6 = [
        ("1. Localizar a Galeria de Estilos:", " Na guia Página Inicial, localize a caixa de Estilos."),
        ("2. Modificar o Estilo:", " Clique com o BOTÃO DIREITO sobre o estilo 'Normal' -> selecione 'Modificar...'."),
        ("3. Fonte e Alinhamento:", " Escolha Arial 12 pt e clique no botão 'Justificado'."),
        ("4. Caixa Parágrafo:", " Clique no botão inferior esquerdo 'Formatar' -> 'Parágrafo...'."),
        ("5. Recuo e Espaçamento:", " Em 'Especial', selecione 'Primeira linha' em 1,25 cm | Em 'Espaçamento entre linhas', selecione '1,5 linha' | Em 'Depois', digite 6 pt."),
        ("6. Salvar:", " Clique em OK duas vezes. Pronto! Todo o texto padrão agora obedece à ABNT.")
    ]
    for i, (title, desc) in enumerate(steps6):
        p = tf6.paragraphs[0] if i == 0 else tf6.add_paragraph()
        p.space_before = Pt(6) if i > 0 else Pt(0)
        r1 = p.add_run()
        r1.text = f"{title} "
        r1.font.bold = True
        r1.font.size = Pt(11.5)
        r1.font.color.rgb = NAVY
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = DARK_GRAY

    # Slide 7: Passo a Passo - Título 1 e 2
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "Passo a Passo — Configurando Título 1 e Título 2")
    add_footer(slide7, 7, total_slides)
    
    cards_s7 = [
        ("CONFIGURAÇÃO DE TÍTULO 1", "Capítulos Principais (ex: 1. INTRODUÇÃO)", [
            "Botão direito sobre 'Título 1' -> Modificar...",
            "Fonte: Arial 14 pt | Negrito | Cor: Preto/Automático.",
            "Formatar -> Parágrafo -> Antes: 12 pt | Depois: 6 pt.",
            "Recuo Especial: Nenhum | Entrelinhas: 1,5 linha.",
            "Atalho rápido de aplicação no texto: Ctrl + Alt + 1."
        ]),
        ("CONFIGURAÇÃO DE TÍTULO 2", "Subseções (ex: 1.1 Histórico de Voo)", [
            "Botão direito sobre 'Título 2' -> Modificar...",
            "Fonte: Arial 12 pt | Negrito | Cor: Preto/Automático.",
            "Formatar -> Parágrafo -> Antes: 6 pt | Depois: 6 pt.",
            "Recuo Especial: Nenhum | Entrelinhas: 1,5 linha.",
            "Atalho rápido de aplicação no texto: Ctrl + Alt + 2."
        ])
    ]
    for i, (ctitle, csub, clist) in enumerate(cards_s7):
        left_pos = Inches(0.6 + i * 4.5)
        card_bg = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.15), Inches(4.3), Inches(3.9))
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = LIGHT_BLUE if i == 0 else WHITE
        card_bg.line.color.rgb = BORDER_COLOR
        
        tb = slide7.shapes.add_textbox(left_pos + Inches(0.2), Inches(1.25), Inches(3.9), Inches(3.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ctitle
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY
        p_sub = tf.add_paragraph()
        p_sub.text = csub
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = BLUE_ACCENT
        p_sub.space_before = Pt(2)
        for item in clist:
            pi = tf.add_paragraph()
            pi.text = f"• {item}"
            pi.font.size = Pt(10)
            pi.font.color.rgb = DARK_GRAY
            pi.space_before = Pt(6)

    # Slide 8: O Painel de Navegação
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "O Painel de Navegação (Ctrl + F)")
    add_footer(slide8, 8, total_slides)
    
    tb8 = slide8.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(4.0))
    tf8 = tb8.text_frame
    tf8.word_wrap = True
    
    nav_topics = [
        ("Como abrir o Painel:", " Pressione o atalho Ctrl + F e clique na aba 'Títulos' (ou acesse a guia Exibir -> marque 'Painel de Navegação')."),
        ("Árvore de Capítulos Interativa:", " O Word exibe todo o esqueleto do documento em tempo real."),
        ("Navegação Instantânea:", " Clicar em qualquer título salta imediatamente para a página correspondente no texto."),
        ("Diagnóstico de Estilos:", " Se um título NÃO aparecer no painel, significa que você esqueceu de aplicar o estilo Título 1 ou Título 2!"),
        ("Reorganização Rápida:", " É possível arrastar capítulos inteiros para cima ou para baixo direto pelo painel de navegação!")
    ]
    for i, (title, desc) in enumerate(nav_topics):
        p = tf8.paragraphs[0] if i == 0 else tf8.add_paragraph()
        p.space_before = Pt(7) if i > 0 else Pt(0)
        r1 = p.add_run()
        r1.text = f"•  {title} "
        r1.font.bold = True
        r1.font.size = Pt(11.5)
        r1.font.color.rgb = NAVY
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = DARK_GRAY

    # Slide 9: Atividade Prática em Laboratório
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "Atividade Prática em Laboratório da Aula 02")
    add_footer(slide9, 9, total_slides)
    
    tb9 = slide9.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(4.0))
    tf9 = tb9.text_frame
    tf9.word_wrap = True
    
    act9 = [
        ("1. Criar o Documento:", " Inicie um novo arquivo e configure Tamanho A4 e Margens (3-3-2-2)."),
        ("2. Configurar os Estilos:", " Configure os estilos Normal, Título 1 e Título 2 conforme os parâmetros ABNT."),
        ("3. Inserir o Texto Técnico:", " Cole as seções do texto sobre História da Aviação Civil."),
        ("4. Aplicar os Estilos:", " Aplique Título 1 nas seções principais, Título 2 nas subseções e Normal nos parágrafos."),
        ("5. Salvar o Arquivo Base:", " Salve como 'Relatorio_ABNT_Aula02.docx' no seu OneDrive institucional."),
        ("Continuação na Aula 03:", " Guarde este arquivo! Na Aula 03 continuaremos inserindo Capa, Quebras de Seção, Numeração ABNT e o Sumário Automático!")
    ]
    for i, (title, desc) in enumerate(act9):
        p = tf9.paragraphs[0] if i == 0 else tf9.add_paragraph()
        p.space_before = Pt(6) if i > 0 else Pt(0)
        r1 = p.add_run()
        r1.text = f"{title} "
        r1.font.bold = True
        r1.font.size = Pt(11.5)
        r1.font.color.rgb = NAVY if i < 5 else RGBColor(0, 120, 50)
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = DARK_GRAY

    # Slide 10: Atalhos
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "Dicas de Produtividade & Atalhos no Word")
    add_footer(slide10, 10, total_slides)
    
    shortcuts2 = [
        ("Ctrl + Alt + 1", "Aplica o estilo Título 1 ao parágrafo atual"),
        ("Ctrl + Alt + 2", "Aplica o estilo Título 2 ao parágrafo atual"),
        ("Ctrl + Shift + N", "Aplica o estilo Normal (limpa formatações avulsas)"),
        ("Ctrl + J", "Aplica o alinhamento Justificado"),
        ("Ctrl + F", "Abre o Painel de Navegação e Busca"),
        ("Ctrl + B ou Ctrl + S", "Salva o documento no disco ou nuvem")
    ]
    
    table_shape = slide10.shapes.add_table(7, 2, Inches(0.8), Inches(1.2), Inches(8.4), Inches(3.6))
    table = table_shape.table
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(5.2)
    
    cell_00 = table.cell(0, 0)
    cell_00.text = "ATALHO DE TECLADO"
    cell_00.fill.solid()
    cell_00.fill.fore_color.rgb = NAVY
    cell_00.text_frame.paragraphs[0].font.bold = True
    cell_00.text_frame.paragraphs[0].font.color.rgb = WHITE
    cell_00.text_frame.paragraphs[0].font.size = Pt(11)
    
    cell_01 = table.cell(0, 1)
    cell_01.text = "AÇÃO NO MS WORD"
    cell_01.fill.solid()
    cell_01.fill.fore_color.rgb = NAVY
    cell_01.text_frame.paragraphs[0].font.bold = True
    cell_01.text_frame.paragraphs[0].font.color.rgb = WHITE
    cell_01.text_frame.paragraphs[0].font.size = Pt(11)
    
    for row_idx, (sc, desc) in enumerate(shortcuts2, start=1):
        c0 = table.cell(row_idx, 0)
        c0.text = sc
        c0.fill.solid()
        c0.fill.fore_color.rgb = LIGHT_BLUE if row_idx % 2 == 1 else WHITE
        p = c0.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(10.5)
        p.font.color.rgb = NAVY
        
        c1 = table.cell(row_idx, 1)
        c1.text = desc
        c1.fill.solid()
        c1.fill.fore_color.rgb = LIGHT_BLUE if row_idx % 2 == 1 else WHITE
        p = c1.text_frame.paragraphs[0]
        p.font.size = Pt(10.5)
        p.font.color.rgb = DARK_GRAY

    output_path = r'c:\projetos\Material\Informática Aplicada a Aeronáutica\Slides\Slides_Aula_02_Word_Tecnico_e_ABNT.pptx'
    prs.save(output_path)
    print(f"Presentation saved successfully: {output_path}")

# ==============================================================================
# MAIN EXECUTION
# ==============================================================================
if __name__ == '__main__':
    from build_aulas_3_4_pptx import build_aula3, build_aula4
    build_aula2()
    build_aula3()
    build_aula4()
