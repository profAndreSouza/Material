import os
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    pptx_path = r'c:\projetos\Material\Informática Aplicada a Aeronáutica\Slides\Slides_Aula_01_Hardware_e_Estacoes_de_Trabalho.pptx'
    img_dir = r'c:\projetos\Material\Informática Aplicada a Aeronáutica\Slides\images'
    
    prs = pptx.Presentation()
    prs.slide_width = Inches(10.0)
    prs.slide_height = Inches(5.625) # 16:9 ratio
    
    blank_layout = prs.slide_layouts[6] # Blank layout
    
    # Colors
    NAVY = RGBColor(15, 32, 67)
    BLUE_ACCENT = RGBColor(0, 114, 206)
    DARK_GRAY = RGBColor(40, 40, 40)
    LIGHT_BG = RGBColor(245, 247, 250)
    WHITE = RGBColor(255, 255, 255)
    GRAY_TEXT = RGBColor(90, 90, 90)
    LIGHT_BLUE = RGBColor(230, 240, 250)
    CARD_BG = RGBColor(255, 255, 255)
    BORDER_COLOR = RGBColor(210, 220, 230)
    
    def add_header(slide, title_text, category_text="INF-117 — INFORMÁTICA APLICADA A AERONÁUTICA"):
        # Header banner shape
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.9))
        header.fill.solid()
        header.fill.fore_color.rgb = NAVY
        header.line.fill.background()
        
        # Header text
        tf = header.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.5)
        tf.margin_top = Inches(0.12)
        
        p0 = tf.paragraphs[0]
        p0.text = category_text.upper()
        p0.font.size = Pt(9)
        p0.font.bold = True
        p0.font.color.rgb = BLUE_ACCENT
        p0.font.name = "Arial"
        
        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = WHITE
        p1.font.name = "Arial"

    def add_footer(slide, current_slide, total_slides=12):
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.25), Inches(9.0), Inches(0.3))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Fatec Sorocaba — CST em Manutenção de Aeronaves  |  Slide {current_slide} de {total_slides}"
        p.font.size = Pt(9)
        p.font.color.rgb = GRAY_TEXT
        p.font.name = "Arial"

    # ==================== SLIDE 1: CAPA ====================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = NAVY
    bg1.line.fill.background()
    
    # Left Content Box
    tb1 = slide1.shapes.add_textbox(Inches(0.6), Inches(0.8), Inches(4.8), Inches(4.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "FATEC SOROCABA — CURSO SUPERIOR DE TECNOLOGIA EM MANUTENÇÃO DE AERONAVES"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.font.name = "Arial"
    
    p = tf1.add_paragraph()
    p.text = "INF-117 — Informática Aplicada a Aeronáutica"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(200, 225, 255)
    p.font.name = "Arial"
    p.space_before = Pt(10)
    
    p = tf1.add_paragraph()
    p.text = "Hardware, Processadores & Estações de Trabalho"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"
    p.space_before = Pt(14)
    
    p = tf1.add_paragraph()
    p.text = "Fundamentos da Arquitetura Computacional, Especificação MRO/Engenharia e Tópico Especial em Computação Quântica"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(220, 230, 245)
    p.font.name = "Arial"
    p.space_before = Pt(10)
    
    p = tf1.add_paragraph()
    p.text = "Prof. André Souza  |  Aula 01"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.font.name = "Arial"
    p.space_before = Pt(24)
    
    # Right Image
    cover_img_path = os.path.join(img_dir, 'cover.jpg')
    if os.path.exists(cover_img_path):
        slide1.shapes.add_picture(cover_img_path, Inches(5.5), Inches(0.8), Inches(4.0), Inches(4.0))

    # ==================== SLIDE 2: OBJETIVOS ====================
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "Objetivos de Aprendizagem")
    add_footer(slide2, 2)
    
    tb2 = slide2.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(4.0))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    
    objs = [
        ("Competências de Hardware:", " Compreender a atuação da CPU, Memória RAM, SSD NVMe e GPU em sistemas computacionais aeronáuticos."),
        ("Arquitetura de Armazenamento e Leitura:", " Avaliar o impacto da velocidade de disco na abertura de manuais técnicos pesados (AMM, IPC, SRM com 3.000+ páginas em PDF) e arquivos CAD 3D."),
        ("Dimensionamento Técnico por Perfil:", " Especificar e justificar computadores adequados para dois ambientes distintos: Hangar/Oficina MRO vs. Engenharia/PCM."),
        ("Diagnóstico Prático no Sistema Operacional:", " Utilizar ferramentas do Windows 11 para identificar gargalos de hardware em tempo real."),
        ("Visão de Fronteira Tecnológica (Computação Quântica):", " Compreender conceitos básicos de qubits e superposição quântica aplicados ao setor aeroespacial.")
    ]
    
    for i, (title, desc) in enumerate(objs):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_before = Pt(10) if i > 0 else Pt(0)
        
        run1 = p.add_run()
        run1.text = f"•  {title}"
        run1.font.bold = True
        run1.font.size = Pt(13)
        run1.font.color.rgb = NAVY
        run1.font.name = "Arial"
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.size = Pt(13)
        run2.font.color.rgb = DARK_GRAY
        run2.font.name = "Arial"

    # ==================== SLIDE 3: HARDWARE VS SOFTWARE ====================
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "Fundamentos — Hardware vs. Software na Aviação")
    add_footer(slide3, 3)
    
    # 3 Cards Layout
    cards = [
        ("HARDWARE AERONÁUTICO", "A infraestrutura física", [
            "Processadores (CPU) de alto desempenho",
            "Memórias RAM DDR4/DDR5 de alta frequência",
            "Armazenamento SSD NVMe PCIe ultra-rápido",
            "Placas de Vídeo (GPU) dedicadas para 3D/CAD",
            "Terminais industriais e monitores antirreflexo"
        ]),
        ("SOFTWARE DE SISTEMA", "A camada de gerenciamento", [
            "Sistemas Operacionais (Windows 11 / Linux)",
            "Gerenciamento de processos e alocação de RAM",
            "Drivers homologados e certificados para engenharia",
            "Segurança de dados e conectividade de rede MRO",
            "Gerenciamento de periféricos de teste"
        ]),
        ("SOFTWARE DE APLICAÇÃO", "As ferramentas de trabalho", [
            "MS Excel (Modelos de $CG$, PERT/CPM e Confiabilidade)",
            "MS Word (Relatórios Técnicos e OS padronizadas)",
            "Softwares CAD (AutoCAD, SolidWorks, CATIA 3D)",
            "Leitores de Manuais (PDF Vetorial AMM/IPC/SRM)",
            "Sistemas ERP/MRO (TRAX, AMOS, SAP Aviation)"
        ])
    ]
    
    for i, (ctitle, csub, clist) in enumerate(cards):
        left_pos = Inches(0.6 + i * 3.0)
        card_bg = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.2), Inches(2.8), Inches(3.8))
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = LIGHT_BLUE if i == 0 else (CARD_BG if i == 1 else LIGHT_BLUE)
        card_bg.line.color.rgb = BORDER_COLOR
        
        tb = slide3.shapes.add_textbox(left_pos + Inches(0.15), Inches(1.3), Inches(2.5), Inches(3.6))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = ctitle
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY
        
        p_sub = tf.add_paragraph()
        p_sub.text = csub
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = BLUE_ACCENT
        p_sub.space_before = Pt(2)
        
        for item in clist:
            pi = tf.add_paragraph()
            pi.text = f"• {item}"
            pi.font.size = Pt(10)
            pi.font.color.rgb = DARK_GRAY
            pi.space_before = Pt(6)

    # ==================== SLIDE 4: CPU ====================
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "O Processador (CPU): O Cérebro do Sistema")
    add_footer(slide4, 4)
    
    tb4 = slide4.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(4.6), Inches(4.0))
    tf4 = tb4.text_frame
    tf4.word_wrap = True
    
    cpu_topics = [
        ("Arquitetura Interna:", " Núcleos (Cores) executam tarefas independentes; Threads simulam núcleos virtuais para maior dinamismo."),
        ("Frequência de Clock (GHz):", " Define a velocidade de ciclos por segundo. Clocks elevados (>= 4.7 GHz) aceleram desenhos e cálculos individuais."),
        ("Memória Cache (L1/L2/L3):", " Memória ultrarrápida dentro do chip da CPU que armazena instruções imediatas para evitar gargalos."),
        ("Impacto Direto no MS Excel:", " O motor de cálculo do Excel utiliza múltiplos núcleos em paralelo para reavaliar planilhas com milhares de fórmulas matriciais e redes PERT/CPM."),
        ("Recomendação Técnica:", " 6 a 8 núcleos (Intel Core i5/i7 ou AMD Ryzen 5/7).")
    ]
    
    for i, (title, desc) in enumerate(cpu_topics):
        p = tf4.paragraphs[0] if i == 0 else tf4.add_paragraph()
        p.space_before = Pt(8) if i > 0 else Pt(0)
        
        r1 = p.add_run()
        r1.text = title
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = NAVY
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(12)
        r2.font.color.rgb = DARK_GRAY
        
    cpu_img_path = os.path.join(img_dir, 'cpu.jpg')
    if os.path.exists(cpu_img_path):
        slide4.shapes.add_picture(cpu_img_path, Inches(5.4), Inches(1.2), Inches(4.0), Inches(3.8))

    # ==================== SLIDE 5: RAM & SSD NVME ====================
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "Memória RAM e Armazenamento SSD NVMe vs. HDD")
    add_footer(slide5, 5)
    
    tb5 = slide5.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(4.6), Inches(4.0))
    tf5 = tb5.text_frame
    tf5.word_wrap = True
    
    ram_topics = [
        ("Memória RAM (Volátil):", " Armazena temporariamente os programas em uso (Excel, CAD, Manuais). Sem RAM suficiente, o Windows usa o disco (*Paging File*), gerando lentidão crítica."),
        ("Capacidade Necessária:", " 16 GB (Hangar/MRO) a 32 GB / 64 GB DDR5 (Engenharia/PCM)."),
        ("SSD NVMe PCIe 4.0 (Persistente):", " Taxas de leitura de 5.000 a 7.000 MB/s. Chips Flash 100% imunes a vibrações de ferramentas pneumáticas em hangares."),
        ("HDD Mecânico (Obsoleto em MRO):", " Apenas ~150 MB/s. Possui agulhas e pratos magnéticos móveis extremamente vulneráveis a impactos e quedas.")
    ]
    
    for i, (title, desc) in enumerate(ram_topics):
        p = tf5.paragraphs[0] if i == 0 else tf5.add_paragraph()
        p.space_before = Pt(8) if i > 0 else Pt(0)
        
        r1 = p.add_run()
        r1.text = title
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = NAVY
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(12)
        r2.font.color.rgb = DARK_GRAY
        
    storage_img_path = os.path.join(img_dir, 'storage.jpg')
    if os.path.exists(storage_img_path):
        slide5.shapes.add_picture(storage_img_path, Inches(5.4), Inches(1.2), Inches(4.0), Inches(3.8))

    # ==================== SLIDE 6: GPU & PERIFÉRICOS ====================
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "Processamento Gráfico (GPU) e Periféricos Ergonomicos")
    add_footer(slide6, 6)
    
    box_gpu = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.2), Inches(4.2), Inches(3.8))
    box_gpu.fill.solid()
    box_gpu.fill.fore_color.rgb = CARD_BG
    box_gpu.line.color.rgb = BORDER_COLOR
    
    tb_gpu = slide6.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(4.0), Inches(3.6))
    tf_gpu = tb_gpu.text_frame
    tf_gpu.word_wrap = True
    
    p = tf_gpu.paragraphs[0]
    p.text = "PLACA DE VÍDEO (GPU)"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = NAVY
    
    items_gpu = [
        ("VRAM Dedicada (GDDR6):", " Permite rotação fluida de modelos CAD 3D de turbinas e fuselagens."),
        ("GPU Integrada vs. Dedicada:", " Vídeo integrado (Intel Iris Xe) atende diagramas 2D; GPU dedicada (NVIDIA RTX) é necessária para engenharia 3D."),
        ("Liberação do Processador:", " A GPU realiza os cálculos geométricos de renderização, liberando a CPU para a lógica de programas.")
    ]
    for t, d in items_gpu:
        pi = tf_gpu.add_paragraph()
        pi.space_before = Pt(8)
        r1 = pi.add_run()
        r1.text = f"• {t}"
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = NAVY
        r2 = pi.add_run()
        r2.text = f" {d}"
        r2.font.size = Pt(11)
        r2.font.color.rgb = DARK_GRAY

    box_per = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.2), Inches(4.2), Inches(3.8))
    box_per.fill.solid()
    box_per.fill.fore_color.rgb = CARD_BG
    box_per.line.color.rgb = BORDER_COLOR
    
    tb_per = slide6.shapes.add_textbox(Inches(5.3), Inches(1.3), Inches(4.0), Inches(3.6))
    tf_per = tb_per.text_frame
    tf_per.word_wrap = True
    
    p = tf_per.paragraphs[0]
    p.text = "PERIFÉRICOS & ERGONOMIA"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = NAVY
    
    items_per = [
        ("Monitores Duplos (QHD):", " Permitem comparar o manual estrutural (SRM) em uma tela enquanto elabora a Ordem de Serviço na outra."),
        ("Tratamento Antirreflexo (*Matte*):", " Essencial contra a forte iluminação industrial dos galpões de hangar."),
        ("Teclado Numérico Dedicado:", " Acelera a digitação constante de códigos e Part Numbers de peças aeronáuticas."),
        ("Conectividade I/O:", " Portas USB 3.2 e Ethernet Gigabit para sincronização com instrumentos de bancada.")
    ]
    for t, d in items_per:
        pi = tf_per.add_paragraph()
        pi.space_before = Pt(8)
        r1 = pi.add_run()
        r1.text = f"• {t}"
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = NAVY
        r2 = pi.add_run()
        r2.text = f" {d}"
        r2.font.size = Pt(11)
        r2.font.color.rgb = DARK_GRAY

    # ==================== SLIDE 7: HANGAR MRO ====================
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "Especificação 1 — Estação para Hangar & Oficina MRO")
    add_footer(slide7, 7)
    
    tb7 = slide7.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(4.6), Inches(4.0))
    tf7 = tb7.text_frame
    tf7.word_wrap = True
    
    hangar_specs = [
        ("CPU:", " Intel Core i5 / AMD Ryzen 5 (6 Cores / 12 Threads)"),
        ("RAM:", " 16 GB DDR4 3200 MHz"),
        ("Armazenamento:", " SSD NVMe M.2 512 GB PCIe 3.0/4.0"),
        ("GPU:", " Vídeo Integrado (Intel Iris Xe / AMD Radeon Vega)"),
        ("Monitor:", " 24\" Full HD com Tratamento Antirreflexo (*Matte*)"),
        ("Justificativa Técnica:", " Manuais AMM ocupam 2 a 4 GB de RAM ao renderizar PDF vetorial. 16 GB garante alternância fluida com o sistema MRO. O SSD NVMe é resistente a choques mecânicos e vibrações do hangar, carregando o sistema em 5s.")
    ]
    
    for i, (title, desc) in enumerate(hangar_specs):
        p = tf7.paragraphs[0] if i == 0 else tf7.add_paragraph()
        p.space_before = Pt(6) if i > 0 else Pt(0)
        
        r1 = p.add_run()
        r1.text = f"• {title}"
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = NAVY
        
        r2 = p.add_run()
        r2.text = f" {desc}"
        r2.font.size = Pt(11)
        r2.font.color.rgb = DARK_GRAY
        
    hangar_img_path = os.path.join(img_dir, 'hangar.jpg')
    if os.path.exists(hangar_img_path):
        slide7.shapes.add_picture(hangar_img_path, Inches(5.4), Inches(1.2), Inches(4.0), Inches(3.8))

    # ==================== SLIDE 8: ENGENHARIA PCM ====================
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "Especificação 2 — Workstation de Engenharia & PCM")
    add_footer(slide8, 8)
    
    box_eng = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.2), Inches(8.8), Inches(3.8))
    box_eng.fill.solid()
    box_eng.fill.fore_color.rgb = LIGHT_BLUE
    box_eng.line.color.rgb = BORDER_COLOR
    
    tb8 = slide8.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(3.6))
    tf8 = tb8.text_frame
    tf8.word_wrap = True
    
    eng_specs = [
        ("Processador (CPU):", " Intel Core i7/i9 ou Ryzen 7/9 (8 a 16 Cores, Clock Boost >= 4.7 GHz)", " Alto clock acelerado para CAD paramétrico e múltiplos núcleos para cálculo matricial do Excel."),
        ("Memória RAM:", " 32 GB a 64 GB DDR5 5600 MHz (Dual Channel)", " Suporta bases de confiabilidade com 500k+ registros e montagens 3D de mais de 20 GB em RAM."),
        ("Armazenamento:", " SSD NVMe 1 TB a 2 TB PCIe 4.0 (Leitura >= 7.000 MB/s)", " Alta taxa de IOPS reduz o tempo de abertura de projetos CAD complexos de minutos para segundos."),
        ("Placa de Vídeo (GPU):", " NVIDIA RTX Profissional Dedicada (8 a 12 GB VRAM GDDR6)", " Renderização 3D de turbinas em tempo real com drivers de precisão homologados."),
        ("Telas:", " 2 Monitores de 27\" IPS QHD (2560x1440) em Braço Articulado", " 77% mais área visual que 1080p, permitindo exibir manual SRM em uma tela e CAD/Excel na outra.")
    ]
    
    for i, (comp, spec, just) in enumerate(eng_specs):
        p = tf8.paragraphs[0] if i == 0 else tf8.add_paragraph()
        p.space_before = Pt(6) if i > 0 else Pt(0)
        
        r1 = p.add_run()
        r1.text = f"• {comp}"
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = NAVY
        
        r2 = p.add_run()
        r2.text = f" {spec}"
        r2.font.size = Pt(11)
        r2.font.bold = True
        r2.font.color.rgb = BLUE_ACCENT
        
        r3 = p.add_run()
        r3.text = f" — {just}"
        r3.font.size = Pt(10)
        r3.font.color.rgb = DARK_GRAY

    # ==================== SLIDE 9: COMPUTAÇÃO QUÂNTICA ====================
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "Tópico Especial — Computação Quântica na Aviação (Pergunta dos Alunos)")
    add_footer(slide9, 9)
    
    tb9 = slide9.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(4.6), Inches(4.0))
    tf9 = tb9.text_frame
    tf9.word_wrap = True
    
    quantum_topics = [
        ("O que é o Qubit?", " Ao contrário do bit clássico (0 ou 1), o *qubit* utiliza princípios quânticos como **Superposição** (0 e 1 simultaneamente) e **Emaranhamento**."),
        ("Aplicações na Aviação:", " Otimização combinatória extrema de malhas aéreas e rotas em tempo real considerando ventos, tráfego e combustível."),
        ("Ciência de Materiais Aeronáuticos:", " Simulação molecular de novas ligas metálicas mais leves e compósitos resistentes a altíssimas temperaturas."),
        ("Manutenção Preditiva Global:", " Alocação ideal de frotas e cronogramas de revisão (*Fleet Scheduling & Maintenance Alignment*)."),
        ("Quântico vs. Clássico:", " A computação quântica **não substitui** computadores pessoais ou workstations de hangar; ela atua como supercomputador acelerador em nuvem para cálculos massivos.")
    ]
    
    for i, (title, desc) in enumerate(quantum_topics):
        p = tf9.paragraphs[0] if i == 0 else tf9.add_paragraph()
        p.space_before = Pt(6) if i > 0 else Pt(0)
        
        r1 = p.add_run()
        r1.text = f"• {title}"
        r1.font.bold = True
        r1.font.size = Pt(11)
        r1.font.color.rgb = NAVY
        
        r2 = p.add_run()
        r2.text = f" {desc}"
        r2.font.size = Pt(11)
        r2.font.color.rgb = DARK_GRAY
        
    quantum_img_path = os.path.join(img_dir, 'quantum.jpg')
    if os.path.exists(quantum_img_path):
        slide9.shapes.add_picture(quantum_img_path, Inches(5.4), Inches(1.2), Inches(4.0), Inches(3.8))

    # ==================== SLIDE 10: DIAGNÓSTICO WINDOWS ====================
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "Prática em Laboratório — Diagnóstico no Windows 11")
    add_footer(slide10, 10)
    
    diag_steps = [
        ("1. Informações do Sistema (`msinfo32` / Settings)", "Verifique o modelo exato do processador (CPU), clock base e quantidade de Memória RAM instalada no computador da bancada.", "Atalho: Win + Pause/Break ou Configurações -> Sistema -> Sobre"),
        ("2. Gerenciador de Tarefas (`Ctrl + Shift + Esc`)", "Abra a aba Desempenho para monitorar a utilização da CPU em tempo real e se a memória está saturada.", "Identifique o tipo de disco (SSD NVMe vs HDD) e a taxa de uso em %"),
        ("3. Identificação de Gargalos Computacionais", "Avalie se a taxa de ocupação da RAM supera 80% e se o disco atinge 100% de uso constante durante a abertura de arquivos grandes.", "Gargalo no disco indica necessidade imediata de upgrade para SSD NVMe")
    ]
    
    for i, (title, desc, tip) in enumerate(diag_steps):
        top_pos = Inches(1.2 + i * 1.3)
        box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), top_pos, Inches(8.8), Inches(1.15))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = BORDER_COLOR
        
        tb = slide10.shapes.add_textbox(Inches(0.8), top_pos + Inches(0.08), Inches(8.4), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY
        
        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(10)
        p1.font.color.rgb = DARK_GRAY
        p1.space_before = Pt(2)
        
        p2 = tf.add_paragraph()
        p2.text = f"💡 Dica Prática: {tip}"
        p2.font.size = Pt(9)
        p2.font.bold = True
        p2.font.color.rgb = BLUE_ACCENT
        p2.space_before = Pt(2)

    # ==================== SLIDE 11: EXERCÍCIO ====================
    slide11 = prs.slides.add_slide(blank_layout)
    add_header(slide11, "Exercício de Fixação Computacional — Parecer Técnico")
    add_footer(slide11, 11)
    
    box_ex = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.2), Inches(8.8), Inches(3.8))
    box_ex.fill.solid()
    box_ex.fill.fore_color.rgb = LIGHT_BLUE
    box_ex.line.color.rgb = BORDER_COLOR
    
    tb_ex = slide11.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(3.6))
    tf_ex = tb_ex.text_frame
    tf_ex.word_wrap = True
    
    p = tf_ex.paragraphs[0]
    p.text = "DESAFIO INDIVIDUAL NO LABORATORIO"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = NAVY
    
    ex_text = [
        "Cenário Prático: Você foi designado para emitir o parecer técnico de aquisição de 2 computadores para a nova base da oficina:",
        "1. Estação de Bancada / Hangar: Lançamento de Ordens de Serviço (OS) e consulta a manuais técnicos AMM/IPC.",
        "2. Workstation de Engenharia & Confiabilidade: Análise de falhas de frota, cálculos em planilhas pesadas e modelos CAD 3D.",
        "Entregável Esperado:",
        "• Monte uma tabela especificada detalhando: CPU (Cores/Clock), RAM (GB/Freq.), Armazenamento (Tipo/Capacidade), GPU e Monitores.",
        "• Apresente a JUSTIFICATIVA TÉCNICA para cada componente selecionado com base nos conceitos aprendidos nesta aula."
    ]
    
    for line in ex_text:
        pi = tf_ex.add_paragraph()
        pi.text = line
        if line.startswith("Cenário") or line.startswith("Entregável"):
            pi.font.bold = True
            pi.font.size = Pt(12)
            pi.font.color.rgb = NAVY
            pi.space_before = Pt(8)
        else:
            pi.font.size = Pt(11)
            pi.font.color.rgb = DARK_GRAY
            pi.space_before = Pt(4)

    # ==================== SLIDE 12: ENCERRAMENTO ====================
    slide12 = prs.slides.add_slide(blank_layout)
    bg12 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg12.fill.solid()
    bg12.fill.fore_color.rgb = NAVY
    bg12.line.fill.background()
    
    tb12 = slide12.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(8.0), Inches(3.5))
    tf12 = tb12.text_frame
    tf12.word_wrap = True
    
    p = tf12.paragraphs[0]
    p.text = "INF-117 — Informática Aplicada a Aeronáutica"
    p.font.size = Pt(14)
    p.font.color.rgb = BLUE_ACCENT
    
    p = tf12.add_paragraph()
    p.text = "Obrigado e Bom Trabalho Prático!"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_before = Pt(12)
    
    p = tf12.add_paragraph()
    p.text = "Próxima Aula (Aula 02): MS Word Técnico para Relatórios Formais — Estilos ABNT, Quebras de Seção, Sumários Automáticos e Tabelas de Engenharia."
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(200, 225, 255)
    p.space_before = Pt(16)
    
    p = tf12.add_paragraph()
    p.text = "Fatec Sorocaba — CST em Manutenção de Aeronaves"
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY_TEXT
    p.space_before = Pt(30)

    prs.save(pptx_path)
    print(f"Presentation saved successfully to {pptx_path}")

if __name__ == "__main__":
    create_presentation()
