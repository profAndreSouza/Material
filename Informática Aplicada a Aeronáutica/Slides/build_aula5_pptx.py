import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Theme Colors
NAVY = RGBColor(15, 32, 67)
BLUE_ACCENT = RGBColor(0, 114, 206)
DARK_GRAY = RGBColor(40, 40, 40)
LIGHT_BG = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)
GRAY_TEXT = RGBColor(90, 90, 90)
CARD_BG = RGBColor(255, 255, 255)
BORDER_COLOR = RGBColor(210, 220, 230)

def add_header(slide, title_text, category_text="INF-117 — INFORMÁTICA APLICADA A AERONÁUTICA"):
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()
    
    tf = header.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.12)
    
    p0 = tf.paragraphs[0]
    p0.text = category_text.upper()
    p0.font.size = Pt(9)
    p0.font.bold = True
    p0.font.color.rgb = BLUE_ACCENT
    p0.font.name = "Arial"
    
    p1 = tf.add_paragraph()
    p1.text = title_text
    p1.font.size = Pt(17)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.font.name = "Arial"

def add_footer(slide, current_slide, total_slides=13):
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.25), Inches(9.0), Inches(0.3))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Fatec Sorocaba — CST em Manutenção de Aeronaves  |  Slide {current_slide} de {total_slides}"
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY_TEXT
    p.font.name = "Arial"

def create_card(slide, left, top, width, height, title, items, bg_color=CARD_BG, border_color=BORDER_COLOR, title_color=NAVY, is_bullet=True):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = border_color
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.16)
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = title_color
    
    for it in items:
        p = tf.add_paragraph()
        p.text = f"• {it}" if is_bullet else it
        p.font.size = Pt(9.5)
        p.font.color.rgb = DARK_GRAY if bg_color != NAVY else RGBColor(220, 235, 250)
        p.space_before = Pt(3)
    return box

def build_aula5():
    prs = pptx.Presentation()
    prs.slide_width = Inches(10.0)
    prs.slide_height = Inches(5.625)
    blank_layout = prs.slide_layouts[6]
    total_slides = 13

    # SLIDE 1: CAPA
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg1.fill.solid(); bg1.fill.fore_color.rgb = NAVY; bg1.line.fill.background()
    tb1 = s1.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8.4), Inches(4.0))
    tf1 = tb1.text_frame; tf1.word_wrap = True
    p = tf1.paragraphs[0]; p.text = "FATEC SOROCABA — CURSO SUPERIOR DE TECNOLOGIA EM MANUTENÇÃO DE AERONAVES"; p.font.size = Pt(10); p.font.bold = True; p.font.color.rgb = BLUE_ACCENT
    p = tf1.add_paragraph(); p.text = "INF-117 — Informática Aplicada a Aeronáutica"; p.font.size = Pt(14); p.font.color.rgb = RGBColor(200, 225, 255); p.space_before = Pt(8)
    p = tf1.add_paragraph(); p.text = "MS Excel I — Primeiros Passos no Excel"; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = WHITE; p.space_before = Pt(14)
    p = tf1.add_paragraph(); p.text = "Interface, Operadores Aritméticos, Precedência com ^ e -, e Análise de Média vs. Mediana com Outliers"; p.font.size = Pt(13); p.font.color.rgb = RGBColor(220, 235, 250)
    p = tf1.add_paragraph(); p.text = "Oficina Prática: Planilha de Custos Operacionais, Física Aeronáutica e Ordens de Serviço"; p.font.size = Pt(11); p.font.color.rgb = BLUE_ACCENT; p.space_before = Pt(20)

    # SLIDE 2: OBJETIVOS
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "Objetivos de Aprendizagem & Contexto Aeronáutico")
    add_footer(s2, 2, total_slides)
    create_card(s2, 0.5, 1.1, 4.3, 3.9, "🎯 Competências Técnicas", [
        "Compreender a matriz de Linhas, Colunas e Células (A1, B4).",
        "A Regra de Ouro do sinal '=' para execução de equações.",
        "Tipagem de dados: Texto, Números, Moeda (R$) e Unidades.",
        "Operadores aritméticos (+, -, *, /, ^) e precedência matemática.",
        "Casos críticos dos operadores ^ e - na física do voo.",
        "Diferença entre Média (=MÉDIA) e Mediana (=MED) com outliers.",
        "Consolidação de tabelas com =SOMA() e =MED()."
    ])
    create_card(s2, 5.1, 1.1, 4.4, 3.9, "✈️ Aplicação na Manutenção", [
        "Controle de Custos de Hangaragem e Estadia de Frotas.",
        "Conciliação de Abastecimento (Litros de AVGAS / JET A-1).",
        "Cálculo de pressão dinâmica e gradiente térmico de subida.",
        "Avaliação correta de custos típicos de Ordens de Serviço (OS) evitando distorções financeiras causadas por outliers."
    ], bg_color=NAVY, title_color=RGBColor(56, 189, 248))

    # SLIDE 3: ANATOMIA E REGRA DO =
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "Anatomia do Excel & A Regra de Ouro do Sinal '='")
    add_footer(s3, 3, total_slides)
    create_card(s3, 0.5, 1.1, 4.3, 3.9, "📐 Endereçamento de Células", [
        "Colunas: Letras maiúsculas (A, B, C... Z, AA...).",
        "Linhas: Números sequenciais (1, 2, 3... 1.048.576).",
        "Célula: Interseção Coluna × Linha (ex: C4 = Coluna C, Linha 4).",
        "Cálculos com referências (ex: =B4*C4) recalculam a planilha automaticamente quando a entrada muda!"
    ])
    create_card(s3, 5.1, 1.1, 4.4, 3.9, "⚠️ Regra de Ouro do Sinal '='", [
        "Toda fórmula DEVE começar com o sinal '='.",
        "100 + 50  →  Texto estático (exibe '100 + 50')",
        "=100 + 50  →  Cálculo aritmético (exibe 150)",
        "=C4 * D4  →  Referência dinâmica calculada",
        "Nunca digite números fixos quando eles já existem na tabela!"
    ], bg_color=RGBColor(255, 251, 235), border_color=RGBColor(253, 230, 138), title_color=RGBColor(180, 83, 9))

    # SLIDE 4: OPERADORES & PRECEDÊNCIA
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "Operadores Aritméticos & Hierarquia de Precedência")
    add_footer(s4, 4, total_slides)
    create_card(s4, 0.5, 1.1, 4.3, 3.9, "➕ Operadores Básicos", [
        "Adição (+): =15 + 5 (Resultado: 20)",
        "Subtração (-): =50 - 20 (Resultado: 30)",
        "Multiplicação (*): =4 * 150 (Resultado: 600)",
        "Divisão (/): =120 / 4 (Resultado: 30)",
        "Exponenciação (^): =2 ^ 3 (Resultado: 8)"
    ])
    create_card(s4, 5.1, 1.1, 4.4, 3.9, "⚡ Armadilha do Sinal de Menos (-)", [
        "1. Parênteses ( ) — Prioridade Absoluta",
        "2. Exponenciação (^) e Negação Unária (-)",
        "3. Multiplicação (*) e Divisão (/)",
        "4. Adição (+) e Subtração (-)",
        "Armadilha no Excel:",
        "=-5^2  →  Calcula (-5)^2 = 25!",
        "Boas Práticas: =-(5^2)  →  Retorna -25 com precisão."
    ], bg_color=RGBColor(240, 247, 255), border_color=RGBColor(186, 230, 253))

    # SLIDE 5: OPERADORES ^ E - NA FÍSICA AERONÁUTICA
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "Operadores '^' e '-' em Fórmulas Aeronáuticas")
    add_footer(s5, 5, total_slides)
    create_card(s5, 0.5, 1.1, 4.3, 3.9, "💨 1. Pressão Dinâmica (q = 0.5 * ρ * V²)", [
        "Fórmula fundamental da sustentação e arrasto:",
        "No Excel: =0.5 * B5 * C5^2",
        "B5: Densidade do ar (kg/m³)",
        "C5: Velocidade verdadeira TAS (m/s)",
        "O operador ^ eleva a velocidade ao quadrado antes de multiplicar pela densidade."
    ])
    create_card(s5, 5.1, 1.1, 4.4, 3.9, "🌡️ 2. Gradiente Térmico ISA (T = T₀ - 0.0065*h)", [
        "Queda de temperatura com a altitude na atmosfera padrão:",
        "No Excel: =15 - (0.0065 * A13)",
        "A13: Altitude em metros (m)",
        "15°C: Temperatura ao nível do mar (T₀)",
        "A multiplicação (0.0065*h) ocorre antes da subtração, garantindo a taxa de 6,5°C/1.000m."
    ], bg_color=NAVY, title_color=RGBColor(56, 189, 248))

    # SLIDE 6: MÉDIA VS MEDIANA
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "Média (=MÉDIA) vs. Mediana (=MED) diante de Outliers")
    add_footer(s6, 6, total_slides)
    create_card(s6, 0.5, 1.1, 4.3, 3.9, "🔧 Estudo de Caso no Hangar (Ordens de Serviço)", [
        "5 OS de Rotina (Troca de óleo, inspeção 50h, magnetos): R$ 1.200 a R$ 2.500.",
        "1 OS Extraordinária (Overhaul de Turbina PT6A): R$ 120.000 (OUTLIER).",
        "MÉDIA (=MÉDIA): R$ 21.500,00 ❌ (Distorcida pelo outlier)",
        "MEDIANA (=MED): R$ 1.900,00 ✅ (Representa o custo rotineiro real)"
    ])
    create_card(s6, 5.1, 1.1, 4.4, 3.9, "📊 Regra de Decisão do Engenheiro", [
        "MÉDIA: Use para orçamento financeiro consolidado, fluxo de caixa e totalizações globais.",
        "MEDIANA: Use para precificação de pacotes de serviços e cálculo de prazos padrão de atendimento (Lead Time) imunes a ruídos estatísticos e eventos raros."
    ], bg_color=RGBColor(240, 253, 244), border_color=RGBColor(187, 247, 208))

    # SLIDE 7: ATIVIDADE 1
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "Atividade 1: Estruturando o Relatório do Hangar")
    add_footer(s7, 7, total_slides)
    create_card(s7, 0.5, 1.1, 4.3, 3.9, "📋 Passo a Passo (Atividade 1)", [
        "1. Renomeie a aba para 'Custos_Hangar'.",
        "2. Título em A1: 'AEROCLUBE DE SOROCABA — CONTROLE DE HANGARAGEM E ABASTECIMENTO' (Mesclar até I1).",
        "3. Cabeçalhos da Linha 3 (9 colunas):",
        "   A3: Prefixo | B3: Modelo da Aeronave",
        "   C3: Diárias Hangar | D3: Valor da Diária (R$)",
        "   E3: Litros Abastecidos (L) | F3: Preço por Litro (R$)",
        "   G3: Subtotal Hangar (R$) | H3: Subtotal Combustível (R$)",
        "   I3: Custo Total (R$)"
    ])
    create_card(s7, 5.1, 1.1, 4.4, 3.9, "🖥️ Organização das Colunas", [
        "Colunas A a F: Dados de Entrada (Brutos).",
        "Colunas G a I: Campos Calculados (Processamento).",
        "Formatação de Cabeçalho: Fundo Azul Marinho (#0F2043) com texto branco em negrito e alinhamento centralizado com quebra automática de texto."
    ], bg_color=NAVY, title_color=RGBColor(56, 189, 248))

    # SLIDE 8: ATIVIDADE 2
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "Atividade 2: Lançamento dos Dados da Frota")
    add_footer(s8, 8, total_slides)
    create_card(s8, 0.5, 1.1, 9.0, 3.9, "✈️ Dados Operacionais da Frota (Linhas 4 a 8)", [
        "Linha 4: PR-ABC | Cessna 172 | 4 diárias | R$ 150,00 | 120 L | R$ 11,50",
        "Linha 5: PT-KRT | Piper Seneca | 10 diárias | R$ 280,00 | 350 L | R$ 11,50",
        "Linha 6: PR-FTE | Beech Baron | 7 diárias | R$ 320,00 | 280 L | R$ 11,50",
        "Linha 7: PP-ZUL | Cirrus SR22 | 3 diárias | R$ 220,00 | 160 L | R$ 11,50",
        "Linha 8: PR-MNT | Cessna 152 | 15 diárias | R$ 120,00 | 90 L | R$ 11,50",
        "Dica: No Excel em português, utilize vírgula (,) para separar as casas decimais (ex: 11,50)."
    ])

    # SLIDE 9: ATIVIDADE 3
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "Atividade 3: Fórmulas Aritméticas & Alça de Preenchimento")
    add_footer(s9, 9, total_slides)
    create_card(s9, 0.5, 1.1, 4.3, 3.9, "⚙️ Equações da Linha 4", [
        "Subtotal Hangar (G4): =C4*D4 (Diárias × Valor Diária)",
        "Subtotal Combustível (H4): =E4*F4 (Litros × Preço/L)",
        "Custo Total (I4): =G4+H4 (Soma dos Subtotais)"
    ])
    create_card(s9, 5.1, 1.1, 4.4, 3.9, "🚀 Alça de Preenchimento", [
        "1. Selecione o bloco G4:I4.",
        "2. Posicione o cursor no quadradinho verde inferior direito.",
        "3. Dê duplo clique para replicar as equações até a linha 8.",
        "O Excel ajusta automaticamente as referências (C5*D5, C6*D6...)."
    ], bg_color=RGBColor(240, 247, 255), border_color=RGBColor(186, 230, 253))

    # SLIDE 10: ATIVIDADE 4
    s10 = prs.slides.add_slide(blank_layout)
    add_header(s10, "Atividade 4: Formatação Moeda (R$) & Função =SOMA()")
    add_footer(s10, 10, total_slides)
    create_card(s10, 0.5, 1.1, 4.3, 3.9, "💰 Formatação Visual e Moeda", [
        "Formato Moeda (R$): Selecione D, F, G, H, I e use Ctrl + Shift + $.",
        "Bordas: Aplique 'Todas as Bordas' de A3 até I8.",
        "Alinhamentos: Centralize Prefixos, Diárias e Litros."
    ])
    create_card(s10, 5.1, 1.1, 4.4, 3.9, "📊 Linha de Totais (Linha 9)", [
        "A9: TOTAL GERAL (Negrito)",
        "G9: =SOMA(G4:G8) (Total Hangar)",
        "H9: =SOMA(H4:H8) (Total Combustível)",
        "I9: =SOMA(I4:I8) (Custo Total Consolidado)",
        "Atalho: Selecione a célula e pressione Alt + = para inserir =SOMA() automaticamente!"
    ], bg_color=RGBColor(240, 253, 244), border_color=RGBColor(187, 247, 208))

    # SLIDE 11: EXERCÍCIO DE FIXAÇÃO
    s11 = prs.slides.add_slide(blank_layout)
    add_header(s11, "Exercício de Fixação: Controle de Consumo de Treinamento")
    add_footer(s11, 11, total_slides)
    create_card(s11, 0.5, 1.1, 4.3, 3.9, "✈️ Aba: Consumo_Treinamento", [
        "Piloto Aluno | Horas Voo (h) | Consumo (L/h) | Consumo Total (L)",
        "Lucas Silveira: 18,5 h | 34,0 L/h",
        "Mariana Duarte: 25,0 h | 36,5 L/h",
        "Carlos Eduardo: 12,0 h | 32,0 L/h",
        "Beatriz Ramos: 30,5 h | 38,0 L/h",
        "Rodrigo Nogueira: 22,0 h | 35,0 L/h",
        "Gabriel Souza: 15,5 h | 33,5 L/h"
    ])
    create_card(s11, 5.1, 1.1, 4.4, 3.9, "🧮 Equações Requeridas", [
        "Consumo por Aluno (D4): =B4*C4 (Horas × Consumo/h)",
        "Total Consumido (D10): =SOMA(D4:D9)",
        "Média por Aluno (D11): =MÉDIA(D4:D9)",
        "Mediana por Aluno (D12): =MED(D4:D9)"
    ], bg_color=RGBColor(240, 247, 255), border_color=RGBColor(186, 230, 253))

    # SLIDE 12: ATALHOS & DICAS
    s12 = prs.slides.add_slide(blank_layout)
    add_header(s12, "Atalhos Rápidos & Boas Práticas em Planilhas Técnicas")
    add_footer(s12, 12, total_slides)
    create_card(s12, 0.5, 1.1, 4.3, 3.9, "⚡ Atalhos Essenciais", [
        "Ctrl + Z / Ctrl + Y: Desfazer / Refazer.",
        "Alt + =: Inserir função =SOMA() automática.",
        "Ctrl + Shift + $: Formatar como Moeda (R$).",
        "Ctrl + Shift + %: Formatar como Porcentagem (%).",
        "=MED(): Cálculo da Mediana central.",
        "Duplo clique na divisa de coluna: Autoajuste de largura."
    ])
    create_card(s12, 5.1, 1.1, 4.4, 3.9, "🛡️ Boas Práticas Aeronáuticas", [
        "Nunca digite totais manualmente: utilize fórmulas dinâmicas.",
        "Indique claramente as unidades nos cabeçalhos ((L), (h), (R$), (Pa)).",
        "Parênteses defensivos: use ( ) em expressões combinadas.",
        "Aplique borda dupla na linha de fechamento contábil."
    ], bg_color=NAVY, title_color=RGBColor(56, 189, 248))

    # SLIDE 13: SÍNTESE & PRÓXIMOS PASSOS
    s13 = prs.slides.add_slide(blank_layout)
    add_header(s13, "Síntese da Aula, Checklist & Próximos Passos")
    add_footer(s13, 13, total_slides)
    create_card(s13, 0.5, 1.1, 4.3, 3.9, "✅ Checklist de Conclusão", [
        "✔ Planilha estruturada com 9 colunas no padrão oficial.",
        "✔ 5 aeronaves cadastradas e formatadas.",
        "✔ Subtotais e Custo Total com fórmulas ativas (* e +).",
        "✔ Física do voo calculada com ^ e - (Pressão Dinâmica e ISA).",
        "✔ Análise crítica de Média vs. Mediana com outliers no hangar.",
        "✔ Exercício de fixação com =SOMA(), =MÉDIA() e =MED()."
    ], bg_color=RGBColor(240, 253, 244), border_color=RGBColor(187, 247, 208))
    create_card(s13, 5.1, 1.1, 4.4, 3.9, "🚀 Próxima Aula — Aula 06: MS Excel II", [
        "Referências Absolutas ($A$1) para fixação de taxas.",
        "Funções Estatísticas: =MÁXIMO(), =MÍNIMO(), =CONT.VALORES().",
        "Gráficos Técnicos de Frota e Desempenho Operacional."
    ], bg_color=NAVY, title_color=RGBColor(56, 189, 248))

    out_pptx = r"c:\projetos\Material\Informática Aplicada a Aeronáutica\Slides\Slides_Aula_05_Excel_Basico_e_Operadores.pptx"
    prs.save(out_pptx)
    print(f"Apresentação PPTX gerada com sucesso em: {out_pptx}")

if __name__ == "__main__":
    build_aula5()
